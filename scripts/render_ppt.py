#!/usr/bin/env python3
"""
render_ppt.py - 고도화된 PPTX 렌더러 v2.1

기능:
- 레이아웃별 렌더링 전략
- 컬럼 기반 레이아웃 지원
- 구조화된 불릿 (레벨, 강조) 지원
- content_blocks 지원 (bullets, table, chart, image, quote, kpi, callout)
- table/kpi/quote/image 최소 렌더링 (placeholder + caption)
- 차트/이미지 플레이스홀더 지원 (chart.data_inline/data_path 시 실제 차트 렌더)
- 디자인 토큰 기반 스타일 적용

개선사항 (v2.1):
- content_blocks 지원: bullets/table/chart/image/quote/kpi 등 블록 타입 렌더 분기
- table/kpi/quote/image에 대한 최소 렌더링 (placeholder+caption) 구현
- layouts.yaml의 render_strategy 참조 (문서화 목적, 실제 렌더링은 코드에서 처리)
"""

import sys
import csv
import json
from pathlib import Path
from typing import Optional, List, Union

import yaml
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement


def load_yaml(path: Path) -> dict:
    """YAML 파일 로드"""
    with path.open("r", encoding="utf-8") as f:
        return yaml.safe_load(f) or {}


def hex_to_rgb(hex_str: str) -> RGBColor:
    """HEX 색상을 RGB로 변환"""
    hex_str = hex_str.strip().lstrip("#")
    if len(hex_str) != 6:
        return RGBColor(0, 0, 0)
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


class DeckRenderer:
    """PPTX 렌더러 클래스 v2.1"""

    def __init__(self, tokens: dict, layouts: dict):
        self.tokens = tokens
        self.layouts = layouts
        self.layout_map = layouts.get("layout_map", {})
        self.prs: Optional[Presentation] = None
        self.asset_base_dir: Optional[Path] = None
        render_options = tokens.get("render_options", {}) if isinstance(tokens.get("render_options", {}), dict) else {}
        self.preserve_template_title_style = bool(render_options.get("preserve_template_title_style", False))
        self.use_body_placeholder_for_bullets = bool(render_options.get("use_body_placeholder_for_bullets", False))
        self.preserve_template_body_style = bool(render_options.get("preserve_template_body_style", False))
        self.force_regular_weight = bool(render_options.get("force_regular_weight", True))
        self.add_chrome_divider = bool(render_options.get("add_chrome_divider", True))
        self.add_layout_chip = bool(render_options.get("add_layout_chip", True))

    def _resolve_bold(self, desired: bool = False) -> bool:
        """사용자 요청 시 전체 폰트 weight를 일반체로 강제"""
        if self.force_regular_weight:
            return False
        return bool(desired)

    def _set_run_font_name(self, run, font_name: str):
        """한글 렌더 호환을 위해 latin/eastAsia/cs 폰트를 함께 지정"""
        name = str(font_name or "Noto Sans KR")
        run.font.name = name
        r = getattr(run, "_r", None)
        if r is None:
            return
        r_pr = r.get_or_add_rPr()
        r_fonts = r_pr.find(qn("a:rFonts"))
        if r_fonts is None:
            r_fonts = OxmlElement("a:rFonts")
            r_pr.append(r_fonts)
        r_fonts.set(qn("a:latin"), name)
        r_fonts.set(qn("a:ea"), name)
        r_fonts.set(qn("a:cs"), name)

    def _resolve_asset_path(self, raw_path: str) -> Optional[Path]:
        """spec 기준 상대경로를 실제 파일 경로로 해석"""
        if not raw_path:
            return None

        resolved = Path(raw_path).expanduser()
        if not resolved.is_absolute() and self.asset_base_dir:
            resolved = (self.asset_base_dir / resolved).resolve()
        return resolved

    @staticmethod
    def _icon_text(icon_name: str) -> str:
        """아이콘 키를 텍스트 심볼로 매핑"""
        icon_map = {
            "check": "✓",
            "checkmark": "✓",
            "arrow": "→",
            "up": "↑",
            "down": "↓",
            "star": "★",
            "dot": "•",
            "risk": "⚠",
            "insight": "◆",
            "idea": "◈",
        }
        key = str(icon_name or "").strip().lower()
        return icon_map.get(key, "")

    @staticmethod
    def _to_number(value):
        """문자열/숫자를 chart용 숫자로 변환"""
        if isinstance(value, (int, float)):
            return value
        try:
            normalized = str(value).strip().replace(",", "")
            if not normalized:
                return 0
            return float(normalized)
        except (TypeError, ValueError):
            return 0

    def _load_chart_data_from_path(self, data_path: str) -> dict:
        """
        data_path(JSON/CSV)에서 차트 데이터 로드.
        지원 포맷:
        - JSON: {"labels":[...], "values":[...]} 또는 {"labels":[...], "series":[{"name":"...", "values":[...]}]}
        - CSV : 첫 열=labels, 이후 열=시리즈 (헤더 필수)
        """
        resolved_path = self._resolve_asset_path(data_path)
        if not resolved_path or not resolved_path.exists():
            return {}

        suffix = resolved_path.suffix.lower()

        try:
            if suffix == ".json":
                payload = json.loads(resolved_path.read_text(encoding="utf-8"))
                if isinstance(payload, dict):
                    labels = payload.get("labels", [])
                    values = payload.get("values", [])
                    series = payload.get("series", [])
                    if labels and (values or series):
                        return {"labels": labels, "values": values, "series": series}

            if suffix == ".csv":
                with resolved_path.open("r", encoding="utf-8-sig", newline="") as f:
                    rows = list(csv.reader(f))

                if len(rows) < 2:
                    return {}

                headers = rows[0]
                if len(headers) < 2:
                    return {}

                labels = [row[0] for row in rows[1:] if row]
                if not labels:
                    return {}

                # 단일 시리즈
                if len(headers) == 2:
                    values = [
                        self._to_number(row[1]) if len(row) > 1 else 0
                        for row in rows[1:]
                    ]
                    return {"labels": labels, "values": values}

                # 다중 시리즈
                series = []
                for col_idx in range(1, len(headers)):
                    series_name = headers[col_idx] or f"Series {col_idx}"
                    series_values = [
                        self._to_number(row[col_idx]) if len(row) > col_idx else 0
                        for row in rows[1:]
                    ]
                    series.append({"name": series_name, "values": series_values})
                return {"labels": labels, "series": series}
        except Exception:
            return {}

        return {}

    def _get_font_config(self, font_key: str) -> dict:
        """
        폰트 설정 가져오기
        개선: name/size_pt 구조 및 레거시 family/size 호환
        """
        fonts = self.tokens.get("fonts", {})
        font_info = fonts.get(font_key, {})

        if isinstance(font_info, dict):
            return {
                "name": font_info.get("name") or font_info.get("family", "Noto Sans KR"),
                "size_pt": font_info.get("size_pt") or font_info.get("size", 12),
                "bold": font_info.get("bold", False)
            }
        elif isinstance(font_info, str):
            return {
                "name": font_info,
                "size_pt": 12,
                "bold": False
            }

        return {
            "name": "Noto Sans KR",
            "size_pt": 12,
            "bold": False
        }

    def _get_color(self, color_key: str) -> str:
        """색상 가져오기"""
        colors = self.tokens.get("colors", {})
        return colors.get(color_key, "1A1A1A")

    def _apply_text_style(
        self,
        text_frame,
        font_key: str,
        color_key: str = "text_dark",
        alignment: PP_ALIGN = PP_ALIGN.LEFT
    ):
        """텍스트 프레임에 스타일 적용"""
        font_cfg = self._get_font_config(font_key)
        color = self._get_color(color_key)

        for p in text_frame.paragraphs:
            p.alignment = alignment
            if not p.runs:
                run = p.add_run()
                run.text = p.text if p.text else ""

            for run in p.runs:
                self._set_run_font_name(run, font_cfg.get("name", "Noto Sans KR"))
                run.font.size = Pt(font_cfg.get("size_pt", 12))
                run.font.bold = self._resolve_bold(font_cfg.get("bold", False))
                run.font.color.rgb = hex_to_rgb(color)

    def _set_title(self, slide, title_text: str):
        """슬라이드 제목 설정"""
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            if not self.preserve_template_title_style:
                self._apply_text_style(
                    slide.shapes.title.text_frame,
                    "title",
                    "text_dark"
                )
        else:
            # 폴백: 텍스트박스 추가
            tx = slide.shapes.add_textbox(Pt(43), Pt(20), Pt(860), Pt(50))
            tf = tx.text_frame
            tf.text = title_text
            self._apply_text_style(tf, "title", "text_dark")

    def _set_body_placeholder_bullets(self, slide, bullets: List[Union[str, dict]]) -> bool:
        """가능하면 BODY/OBJECT placeholder에 불릿을 채워 템플릿 스타일을 최대한 활용"""
        if not bullets:
            return False

        target = None
        for shape in slide.shapes:
            if not getattr(shape, "is_placeholder", False):
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            placeholder_type = str(shape.placeholder_format.type)
            if "BODY" in placeholder_type or "OBJECT" in placeholder_type:
                target = shape
                break

        if not target:
            return False

        tf = target.text_frame
        tf.clear()
        tf.word_wrap = True

        body_font_cfg = self._get_font_config("body")
        color = self._get_color("text_dark")

        for idx, bullet in enumerate(bullets):
            if isinstance(bullet, str):
                text = bullet
                level = 0
                emphasis = "normal"
                icon = ""
            else:
                text = str(bullet.get("text", ""))
                level = bullet.get("level", 0)
                emphasis = bullet.get("emphasis", "normal")
                icon = self._icon_text(bullet.get("icon", ""))

            if icon:
                text = f"{icon} {text}"

            if idx == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = text
            para.level = level
            para.alignment = PP_ALIGN.LEFT

            if self.preserve_template_body_style:
                continue

            for run in para.runs if para.runs else [para.add_run()]:
                if not para.runs:
                    run.text = text
                self._set_run_font_name(run, body_font_cfg.get("name", "Noto Sans KR"))
                base_size = body_font_cfg.get("size_pt", 12)
                if level == 1:
                    run.font.size = Pt(base_size - 1)
                elif level == 2:
                    run.font.size = Pt(base_size - 2)
                else:
                    run.font.size = Pt(base_size)

                if emphasis == "bold":
                    run.font.bold = self._resolve_bold(True)
                elif emphasis == "highlight":
                    run.font.bold = self._resolve_bold(True)
                    run.font.color.rgb = hex_to_rgb(self._get_color("primary_blue"))
                else:
                    run.font.bold = self._resolve_bold(body_font_cfg.get("bold", False))
                    run.font.color.rgb = hex_to_rgb(color)

        return True

    def _add_governing_message(self, slide, msg: str, top_pt: int = 75):
        """거버닝 메시지 추가"""
        tx = slide.shapes.add_textbox(Pt(43), Pt(top_pt), Pt(860), Pt(58))
        tf = tx.text_frame
        tf.text = msg
        tf.word_wrap = True
        self._apply_text_style(tf, "governing", "text_muted")

    def _add_slide_chrome(self, slide, layout_name: str):
        """슬라이드 상단 크롬(구분선/레이아웃 칩) 추가"""
        if not (self.add_chrome_divider or self.add_layout_chip):
            return

        if layout_name in {"cover"}:
            return

        if self.add_chrome_divider:
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Pt(43), Pt(108), Pt(860), Pt(1.4)
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = hex_to_rgb(self._get_color("divider_gray"))
            divider.line.fill.background()

        if self.add_layout_chip and layout_name not in {"section_divider", "thank_you"}:
            chip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Pt(790), Pt(72), Pt(113), Pt(24)
            )
            chip.fill.solid()
            chip.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
            chip.line.fill.background()

            tf = chip.text_frame
            tf.text = str(layout_name).replace("_", " ").upper()
            tf.word_wrap = False
            self._apply_text_style(tf, "footnote", "dark_blue", PP_ALIGN.CENTER)

    def _add_bullets(
        self,
        slide,
        bullets: List[Union[str, dict]],
        left_pt: int = 43,
        top_pt: int = 130,
        width_pt: int = 860,
        height_pt: int = 350,
        color_key: str = "text_dark"
    ):
        """불릿 포인트 추가"""
        if not bullets:
            return

        tx = slide.shapes.add_textbox(
            Pt(left_pt), Pt(top_pt),
            Pt(width_pt), Pt(height_pt)
        )
        tf = tx.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            # 불릿 파싱
            if isinstance(bullet, str):
                text = bullet
                level = 0
                emphasis = "normal"
                icon = ""
            else:
                text = bullet.get("text", "")
                level = bullet.get("level", 0)
                emphasis = bullet.get("emphasis", "normal")
                icon = self._icon_text(bullet.get("icon", ""))

            if icon:
                text = f"{icon} {text}"

            # 첫 번째 불릿은 기존 paragraph 사용
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            p.level = level
            p.alignment = PP_ALIGN.LEFT

            # 스타일 적용
            font_cfg = self._get_font_config("body")
            color = self._get_color(color_key)

            for run in p.runs if p.runs else [p.add_run()]:
                if not p.runs:
                    run.text = text
                self._set_run_font_name(run, font_cfg.get("name", "Noto Sans KR"))

                # 레벨에 따른 폰트 크기 조정
                base_size = font_cfg.get("size_pt", 12)
                if level == 1:
                    run.font.size = Pt(base_size - 1)
                elif level == 2:
                    run.font.size = Pt(base_size - 2)
                else:
                    run.font.size = Pt(base_size)

                # 강조 처리
                if emphasis == "bold":
                    run.font.bold = self._resolve_bold(True)
                elif emphasis == "highlight":
                    run.font.bold = self._resolve_bold(True)
                    run.font.color.rgb = hex_to_rgb(self._get_color("primary_blue"))
                else:
                    run.font.bold = self._resolve_bold(font_cfg.get("bold", False))
                    run.font.color.rgb = hex_to_rgb(color)

    def _add_notes(self, slide, notes: str):
        """발표자 노트 추가"""
        if not notes:
            return
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes

    def _add_footnotes(self, slide, footnotes: List[dict], bottom_pt: int = 490):
        """각주 추가"""
        if not footnotes:
            return

        tx = slide.shapes.add_textbox(Pt(43), Pt(bottom_pt), Pt(860), Pt(30))
        tf = tx.text_frame

        for i, fn in enumerate(footnotes):
            marker = fn.get("marker", f"*{i+1}")
            text = fn.get("text", "")

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"{marker} {text}"

        self._apply_text_style(tf, "footnote", "text_muted")

    def _pick_layout(self, layout_name: str):
        """레이아웃 선택"""
        layout_cfg = self.layout_map.get(layout_name, {"slide_layout_index": 1})
        layout_index = int(layout_cfg.get("slide_layout_index", 1))

        if layout_index < 0 or layout_index >= len(self.prs.slide_layouts):
            # 안전한 폴백
            return self.prs.slide_layouts[1] if len(self.prs.slide_layouts) > 1 else self.prs.slide_layouts[0]

        return self.prs.slide_layouts[layout_index]

    # =========================================================================
    # content_blocks 렌더링 메서드 (v2.1 신규)
    # =========================================================================

    def _estimate_bullet_block_height(self, bullets: List[Union[str, dict]]) -> int:
        """불릿 수/길이를 기반으로 불릿 블록 높이를 추정"""
        if not bullets:
            return 100

        total_lines = 0
        for bullet in bullets:
            text = bullet if isinstance(bullet, str) else str(bullet.get("text", ""))
            # 문장형 불릿 수용을 위해 4줄까지 허용
            lines = max(1, min(4, (len(text) // 38) + 1))
            total_lines += lines

        # 라인당 약 20pt + 상하 여백
        estimated = 32 + (total_lines * 20)
        return max(96, min(estimated, 420))

    def _slide_height_pt(self) -> int:
        """현재 프레젠테이션 슬라이드 높이(pt)"""
        if not self.prs:
            return 540
        return int(self.prs.slide_height / 12700)

    def _fit_block_height(
        self,
        current_top: int,
        desired_height: int,
        min_height: int,
        bottom_margin_pt: int = 8
    ) -> Optional[int]:
        """
        블록 높이를 슬라이드 높이에 맞춰 보정.
        배치 불가능 시 None 반환.
        """
        slide_height_pt = self._slide_height_pt()
        available = slide_height_pt - bottom_margin_pt - current_top
        if available < min_height:
            return None
        return int(max(min_height, min(desired_height, available)))

    def _render_content_blocks(self, slide, content_blocks: List[dict], start_top_pt: int = 130):
        """
        content_blocks 배열 렌더링
        지원 타입: bullets, table, chart, image, quote, kpi, callout, text
        """
        if not content_blocks:
            return

        lane_tops = {
            "main": start_top_pt,
            "left": start_top_pt,
            "right": start_top_pt,
            "sidebar": start_top_pt,
            "top": start_top_pt,
            "bottom": start_top_pt,
        }

        for block in content_blocks:
            block_type = block.get("type", "bullets")
            position = block.get("position", "main")
            lane = position if position in lane_tops else "main"
            current_top = lane_tops[lane]
            custom_top = block.get("top_pt")
            if isinstance(custom_top, (int, float)):
                current_top = int(custom_top)
            current_top = max(110, current_top)
            custom_height = block.get("height_pt")

            # 위치에 따른 좌표 조정
            if isinstance(block.get("left_pt"), (int, float)) and isinstance(block.get("width_pt"), (int, float)):
                left_pt = int(block.get("left_pt"))
                width_pt = int(block.get("width_pt"))
            elif position == "left":
                left_pt, width_pt = 43, 400
            elif position == "right":
                left_pt, width_pt = 480, 400
            elif position == "sidebar":
                left_pt, width_pt = 700, 200
            else:  # main
                left_pt, width_pt = 43, 860

            # 타입별 렌더링
            if block_type == "bullets":
                bullets = block.get("bullets", [])
                block_height = self._estimate_bullet_block_height(bullets)
                if isinstance(custom_height, (int, float)):
                    block_height = max(80, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=76)
                if fitted_height is None:
                    continue
                self._add_bullets(slide, bullets, left_pt, current_top, width_pt, fitted_height)
                lane_tops[lane] = current_top + fitted_height + 18  # 다음 블록 위치

            elif block_type == "table":
                table_def = block.get("table", {})
                rows = table_def.get("rows", [])
                block_height = max(120, min(220, 35 + (len(rows) * 28)))
                if isinstance(custom_height, (int, float)):
                    block_height = max(110, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=110)
                if fitted_height is None:
                    continue
                self._render_table_placeholder(slide, table_def, left_pt, current_top, width_pt)
                lane_tops[lane] = current_top + fitted_height + 18

            elif block_type == "chart":
                chart_def = block.get("chart", {})
                block_height = 230
                if isinstance(custom_height, (int, float)):
                    block_height = max(150, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=150)
                if fitted_height is None:
                    continue
                self._render_chart_placeholder(slide, chart_def, left_pt, current_top, width_pt)
                lane_tops[lane] = current_top + fitted_height + 18

            elif block_type == "image":
                image_def = block.get("image", {})
                block_height = 230
                if isinstance(custom_height, (int, float)):
                    block_height = max(150, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=150)
                if fitted_height is None:
                    continue
                self._render_image_placeholder(slide, image_def, left_pt, current_top, width_pt)
                lane_tops[lane] = current_top + fitted_height + 18

            elif block_type == "quote":
                quote_def = block.get("quote", {})
                block_height = 120
                if isinstance(custom_height, (int, float)):
                    block_height = max(90, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=90)
                if fitted_height is None:
                    continue
                self._render_quote(slide, quote_def, left_pt, current_top, width_pt)
                lane_tops[lane] = current_top + fitted_height + 18

            elif block_type == "kpi":
                kpi_def = block.get("kpi", {})
                block_height = 90
                if isinstance(custom_height, (int, float)):
                    block_height = max(70, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=70)
                if fitted_height is None:
                    continue
                self._render_kpi(slide, kpi_def, left_pt, current_top, width_pt)
                lane_tops[lane] = current_top + fitted_height + 18

            elif block_type == "callout":
                callout_def = block.get("callout", {})
                block_height = 80
                if isinstance(custom_height, (int, float)):
                    block_height = max(60, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=56)
                if fitted_height is None:
                    continue
                self._render_callout(slide, callout_def, left_pt, current_top, width_pt)
                lane_tops[lane] = current_top + fitted_height + 18

            elif block_type == "text":
                text = block.get("text", "")
                block_height = 70
                if isinstance(custom_height, (int, float)):
                    block_height = max(60, int(custom_height))
                fitted_height = self._fit_block_height(current_top, block_height, min_height=54)
                if fitted_height is None:
                    continue
                self._render_text_block(slide, text, left_pt, current_top, width_pt, fitted_height)
                lane_tops[lane] = current_top + fitted_height + 18

            # main lane 블록은 전체 기준선을 함께 이동
            if lane == "main":
                baseline = lane_tops["main"]
                for key in ("left", "right", "sidebar", "top", "bottom"):
                    lane_tops[key] = max(lane_tops[key], baseline)

    def _render_table_placeholder(self, slide, table_def: dict, left_pt: int, top_pt: int, width_pt: int):
        """
        테이블 렌더링 (실제 테이블 또는 플레이스홀더)
        """
        headers = table_def.get("headers", [])
        rows = table_def.get("rows", [])
        style = table_def.get("style", "default")

        if headers and rows:
            # 실제 테이블 생성
            num_rows = len(rows) + 1  # 헤더 포함
            num_cols = len(headers)

            table = slide.shapes.add_table(
                num_rows, num_cols,
                Pt(left_pt), Pt(top_pt),
                Pt(width_pt), Pt(min(30 * num_rows, 200))
            ).table

            # 헤더 설정
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                # 헤더 스타일
                for para in cell.text_frame.paragraphs:
                    para.font.bold = self._resolve_bold(True)
                    para.font.size = Pt(10)
                    for run in para.runs:
                        self._set_run_font_name(run, self._get_font_config("body").get("name", "Noto Sans KR"))

            # 데이터 행 설정
            for row_idx, row in enumerate(rows):
                for col_idx, cell_text in enumerate(row):
                    if col_idx < num_cols:
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_text)
                        for para in cell.text_frame.paragraphs:
                            para.font.size = Pt(9)
                            for run in para.runs:
                                self._set_run_font_name(run, self._get_font_config("body").get("name", "Noto Sans KR"))
        else:
            # 플레이스홀더로 테이블 영역 표시
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(150)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
            shape.line.color.rgb = hex_to_rgb(self._get_color("divider_gray"))

            tf = shape.text_frame
            tf.text = "[TABLE]\n테이블 데이터 필요"
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
            self._apply_text_style(tf, "body", "text_muted", PP_ALIGN.CENTER)

    def _render_chart_placeholder(self, slide, chart_def: dict, left_pt: int, top_pt: int, width_pt: int):
        """
        차트 렌더링 (data_inline/data_path가 있으면 실제 차트, 없으면 플레이스홀더)
        """
        chart_type = chart_def.get("type", "bar_chart")
        title = chart_def.get("title", "")
        caption = chart_def.get("caption", "")
        data_inline = chart_def.get("data_inline", {})
        data_path = chart_def.get("data_path", "")
        rendered = False

        # data_path가 주어지고 data_inline이 비어있으면 파일 로드
        if not data_inline and data_path:
            data_inline = self._load_chart_data_from_path(data_path)

        # 차트 생성 시도
        if data_inline:
            labels = data_inline.get("labels", [])
            values = data_inline.get("values", [])
            series = data_inline.get("series", [])

            if labels and (values or series):
                try:
                    chart_data = ChartData()
                    chart_data.categories = labels

                    if series:
                        for series_item in series:
                            series_name = series_item.get("name", "Series")
                            series_values = [
                                self._to_number(v) for v in series_item.get("values", [])
                            ]
                            chart_data.add_series(series_name, series_values)
                    else:
                        series_name = title if title else "Series"
                        chart_data.add_series(series_name, [self._to_number(v) for v in values])

                    chart_type_map = {
                        "bar_chart": XL_CHART_TYPE.BAR_CLUSTERED,
                        "line_chart": XL_CHART_TYPE.LINE_MARKERS,
                        "pie_chart": XL_CHART_TYPE.PIE,
                        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
                        "scatter": XL_CHART_TYPE.XY_SCATTER,
                    }
                    chart_kind = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

                    chart_shape = slide.shapes.add_chart(
                        chart_kind,
                        Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(200),
                        chart_data
                    )
                    chart = chart_shape.chart
                    chart.has_title = bool(title)
                    if title:
                        chart.chart_title.text_frame.text = title
                    rendered = True
                except Exception:
                    # 차트 생성 실패 시 플레이스홀더로 폴백
                    rendered = False

        if not rendered:
            # 플레이스홀더 박스
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(200)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
            shape.line.color.rgb = hex_to_rgb(self._get_color("divider_gray"))

            # 플레이스홀더 텍스트
            tf = shape.text_frame
            tf.text = f"[CHART: {chart_type}]\n{title}" if title else f"[CHART: {chart_type}]"
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
            self._apply_text_style(tf, "body", "text_muted", PP_ALIGN.CENTER)

        # 캡션
        if caption:
            tx_cap = slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt + 205), Pt(width_pt), Pt(25))
            tx_cap.text_frame.text = caption
            self._apply_text_style(tx_cap.text_frame, "footnote", "text_muted")

    def _render_image_placeholder(self, slide, image_def: dict, left_pt: int, top_pt: int, width_pt: int):
        """
        이미지 플레이스홀더 렌더링
        image_path가 있으면 실제 이미지 삽입 시도, 없으면 플레이스홀더
        """
        image_path = image_def.get("image_path", "")
        title = image_def.get("title", "")
        caption = image_def.get("caption", "")

        resolved_image_path = self._resolve_asset_path(image_path)

        if resolved_image_path and resolved_image_path.exists():
            # 실제 이미지 삽입
            try:
                slide.shapes.add_picture(
                    str(resolved_image_path),
                    Pt(left_pt), Pt(top_pt),
                    width=Pt(width_pt)
                )
            except Exception:
                # 실패 시 플레이스홀더로 폴백
                self._render_image_placeholder_box(slide, title, left_pt, top_pt, width_pt)
        else:
            self._render_image_placeholder_box(slide, title, left_pt, top_pt, width_pt)

        # 캡션
        if caption:
            tx_cap = slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt + 205), Pt(width_pt), Pt(25))
            tx_cap.text_frame.text = caption
            self._apply_text_style(tx_cap.text_frame, "footnote", "text_muted")

    def _render_image_placeholder_box(self, slide, title: str, left_pt: int, top_pt: int, width_pt: int):
        """이미지 플레이스홀더 박스"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(200)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
        shape.line.color.rgb = hex_to_rgb(self._get_color("divider_gray"))

        tf = shape.text_frame
        tf.text = f"[IMAGE]\n{title}" if title else "[IMAGE]\n이미지 경로 필요"
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
        self._apply_text_style(tf, "body", "text_muted", PP_ALIGN.CENTER)

    def _render_quote(self, slide, quote_def: dict, left_pt: int, top_pt: int, width_pt: int):
        """인용문 렌더링"""
        text = quote_def.get("text", "")
        author = quote_def.get("author", "")
        role = quote_def.get("role", "")
        style = quote_def.get("style", "standard")

        # 인용문 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(100)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
        shape.line.fill.background()

        # 인용문 텍스트
        tf = shape.text_frame
        tf.word_wrap = True
        quote_text = f'"{text}"'
        if author:
            attribution = f"— {author}"
            if role:
                attribution += f", {role}"
            quote_text += f"\n{attribution}"

        tf.text = quote_text
        self._apply_text_style(tf, "governing" if style == "large" else "body", "text_dark", PP_ALIGN.CENTER)

    def _render_kpi(self, slide, kpi_def: dict, left_pt: int, top_pt: int, width_pt: int):
        """KPI/지표 렌더링"""
        label = kpi_def.get("label", "KPI")
        value = kpi_def.get("value", "0")
        unit = kpi_def.get("unit", "")
        trend = kpi_def.get("trend", "neutral")
        comparison = kpi_def.get("comparison", "")

        # KPI 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Pt(left_pt), Pt(top_pt), Pt(min(width_pt, 200)), Pt(70)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self._get_color("primary_blue"))
        shape.line.fill.background()

        # KPI 텍스트
        tf = shape.text_frame

        # 값 + 단위
        value_text = f"{value}{unit}"
        if trend == "up":
            value_text = f"↑ {value_text}"
        elif trend == "down":
            value_text = f"↓ {value_text}"

        tf.text = f"{label}\n{value_text}"
        if comparison:
            tf.text += f"\n{comparison}"

        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.color.rgb = hex_to_rgb(self._get_color("background"))
                self._set_run_font_name(run, self._get_font_config("body").get("name", "Noto Sans KR"))
                run.font.size = Pt(12)
                run.font.bold = self._resolve_bold(True)

    def _render_callout(self, slide, callout_def: dict, left_pt: int, top_pt: int, width_pt: int):
        """콜아웃 렌더링"""
        text = callout_def.get("text", "")
        callout_type = callout_def.get("type", "info")
        icon = self._icon_text(callout_def.get("icon", ""))

        # 타입별 색상
        type_colors = {
            "info": "primary_blue",
            "warning": "text_muted",
            "success": "primary_blue",
            "highlight": "dark_blue",
            "key_insight": "dark_blue"
        }
        bg_color = type_colors.get(callout_type, "primary_blue")

        # 콜아웃 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(60)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self._get_color(bg_color))
        shape.line.fill.background()

        tf = shape.text_frame
        tf.text = f"{icon} {text}".strip() if icon else text
        tf.word_wrap = True

        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.color.rgb = hex_to_rgb(self._get_color("background"))
                self._set_run_font_name(run, self._get_font_config("body").get("name", "Noto Sans KR"))
                run.font.size = Pt(11)

    def _render_text_block(self, slide, text: str, left_pt: int, top_pt: int, width_pt: int, height_pt: int = 50):
        """일반 텍스트 블록 렌더링"""
        tx = slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
        tf = tx.text_frame
        tf.text = text
        tf.word_wrap = True
        self._apply_text_style(tf, "body", "text_dark")

    def _extract_chart_visual(self, slide_data: dict) -> Optional[dict]:
        """chart_focus용 차트 정의 추출 (visuals 우선, content_blocks 폴백)"""
        chart_like_types = {"chart", "bar_chart", "line_chart", "pie_chart", "stacked_bar", "scatter"}

        for visual in slide_data.get("visuals", []):
            if not isinstance(visual, dict):
                continue
            v_type = str(visual.get("type", "")).strip().lower()
            if v_type in chart_like_types or visual.get("data_inline") or visual.get("data_path"):
                return visual

        for block in slide_data.get("content_blocks", []):
            if isinstance(block, dict) and block.get("type") == "chart" and isinstance(block.get("chart"), dict):
                return block.get("chart")

        return None

    def _extract_image_visual(self, slide_data: dict) -> Optional[dict]:
        """image_focus용 이미지 정의 추출 (visuals 우선, content_blocks 폴백)"""
        image_like_types = {"image", "photo", "illustration"}

        for visual in slide_data.get("visuals", []):
            if not isinstance(visual, dict):
                continue
            v_type = str(visual.get("type", "")).strip().lower()
            if v_type in image_like_types or visual.get("image_path"):
                return visual

        for block in slide_data.get("content_blocks", []):
            if isinstance(block, dict) and block.get("type") == "image" and isinstance(block.get("image"), dict):
                return block.get("image")

        return None

    # =========================================================================
    # 레이아웃별 렌더링 메서드
    # =========================================================================

    def _render_cover(self, slide, slide_data: dict):
        """커버 슬라이드 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))

        # 부제목 또는 거버닝 메시지를 서브타이틀로
        subtitle = slide_data.get("subtitle") or slide_data.get("governing_message", "")
        if subtitle:
            # 서브타이틀 위치 (커버 중앙 하단)
            tx = slide.shapes.add_textbox(Pt(43), Pt(280), Pt(860), Pt(60))
            tf = tx.text_frame
            tf.text = subtitle
            tf.word_wrap = True
            self._apply_text_style(tf, "governing", "text_muted", PP_ALIGN.CENTER)

        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_section_divider(self, slide, slide_data: dict):
        """섹션 구분 슬라이드 렌더링"""
        title = slide_data.get("title", "")
        governing = slide_data.get("governing_message", "")

        # 중앙 정렬 제목
        tx = slide.shapes.add_textbox(Pt(43), Pt(200), Pt(860), Pt(80))
        tf = tx.text_frame
        tf.text = title
        tf.word_wrap = True
        self._apply_text_style(tf, "title", "primary_blue", PP_ALIGN.CENTER)

        # 부제목
        if governing:
            tx2 = slide.shapes.add_textbox(Pt(43), Pt(290), Pt(860), Pt(50))
            tf2 = tx2.text_frame
            tf2.text = governing
            tf2.word_wrap = True
            self._apply_text_style(tf2, "governing", "text_muted", PP_ALIGN.CENTER)

        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_exec_summary(self, slide, slide_data: dict):
        """Executive Summary 슬라이드 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        # content_blocks + bullets 병행 지원 (밀도 높은 summary 작성)
        if slide_data.get("content_blocks"):
            blocks = list(slide_data.get("content_blocks", []))
            bullets = slide_data.get("bullets", [])
            if bullets and not any(isinstance(b, dict) and b.get("type") == "bullets" for b in blocks):
                blocks.insert(0, {"type": "bullets", "bullets": bullets})
            self._render_content_blocks(slide, blocks)
        else:
            bullets = slide_data.get("bullets", [])
            if not (self.use_body_placeholder_for_bullets and self._set_body_placeholder_bullets(slide, bullets)):
                self._add_bullets(slide, bullets)

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_inline_columns(self, slide, columns: List[dict], start_top_pt: int = 160, body_height: int = 210) -> bool:
        """content 레이아웃에서도 columns 데이터를 실제로 표시"""
        if not columns or len(columns) < 2:
            return False

        if len(columns) >= 3:
            col_positions = [43, 330, 617]
            col_widths = [270, 270, 270]
            max_cols = 3
        else:
            col_positions = [43, 480]
            col_widths = [400, 400]
            max_cols = 2

        for i, col in enumerate(columns[:max_cols]):
            left_pt = col_positions[i]
            width_pt = col_widths[i]
            heading = str(col.get("heading", "")).strip()
            if heading:
                tx_head = slide.shapes.add_textbox(Pt(left_pt), Pt(125), Pt(width_pt), Pt(30))
                tf = tx_head.text_frame
                tf.text = heading
                self._apply_text_style(tf, "governing", "primary_blue")

            if col.get("content_blocks"):
                column_blocks = []
                for block in col.get("content_blocks", []):
                    block_copy = dict(block)
                    block_copy.setdefault("left_pt", left_pt)
                    block_copy.setdefault("width_pt", width_pt)
                    block_copy.setdefault("position", "main")
                    column_blocks.append(block_copy)
                self._render_content_blocks(slide, column_blocks, start_top_pt=start_top_pt)
            else:
                self._add_bullets(
                    slide,
                    col.get("bullets", []),
                    left_pt=left_pt,
                    top_pt=start_top_pt,
                    width_pt=width_pt,
                    height_pt=body_height,
                )

        return True

    def _render_content(self, slide, slide_data: dict):
        """일반 컨텐츠 슬라이드 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        columns = slide_data.get("columns", []) if isinstance(slide_data.get("columns", []), list) else []
        rendered_columns = self._render_inline_columns(slide, columns, start_top_pt=160, body_height=220)

        # content_blocks + bullets 병행 지원 (표/콜아웃 + 핵심 불릿 동시 표현)
        if slide_data.get("content_blocks"):
            blocks = list(slide_data.get("content_blocks", []))
            bullets = slide_data.get("bullets", [])
            if bullets and not any(isinstance(b, dict) and b.get("type") == "bullets" for b in blocks):
                blocks.insert(0, {"type": "bullets", "bullets": bullets})

            if rendered_columns:
                normalized_blocks = []
                for block in blocks:
                    block_copy = dict(block)
                    if not isinstance(block_copy.get("top_pt"), (int, float)):
                        block_copy["top_pt"] = 392
                    normalized_blocks.append(block_copy)
                self._render_content_blocks(slide, normalized_blocks, start_top_pt=392)
            else:
                self._render_content_blocks(slide, blocks)
        else:
            bullets = slide_data.get("bullets", [])
            if rendered_columns:
                if bullets:
                    self._render_content_blocks(
                        slide,
                        [{"type": "bullets", "top_pt": 392, "bullets": bullets}],
                        start_top_pt=392,
                    )
            elif not (self.use_body_placeholder_for_bullets and self._set_body_placeholder_bullets(slide, bullets)):
                self._add_bullets(slide, bullets)

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_two_column(self, slide, slide_data: dict):
        """2컬럼 레이아웃 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        columns = slide_data.get("columns", [])

        if columns and len(columns) >= 2:
            # columns 데이터 사용
            left_col = columns[0]
            right_col = columns[1]

            # 왼쪽 컬럼 헤딩
            tx_left_head = slide.shapes.add_textbox(Pt(43), Pt(125), Pt(400), Pt(30))
            tf = tx_left_head.text_frame
            tf.text = left_col.get("heading", "")
            self._apply_text_style(tf, "governing", "primary_blue")

            # 왼쪽 컬럼: content_blocks 또는 bullets
            if left_col.get("content_blocks"):
                left_blocks = []
                for block in left_col["content_blocks"]:
                    left_block = dict(block)
                    left_block["position"] = "left"
                    left_blocks.append(left_block)
                self._render_content_blocks(slide, left_blocks, start_top_pt=160)
            else:
                self._add_bullets(
                    slide,
                    left_col.get("bullets", []),
                    left_pt=43, top_pt=160, width_pt=400, height_pt=300
                )

            # 오른쪽 컬럼 헤딩
            tx_right_head = slide.shapes.add_textbox(Pt(480), Pt(125), Pt(400), Pt(30))
            tf = tx_right_head.text_frame
            tf.text = right_col.get("heading", "")
            self._apply_text_style(tf, "governing", "primary_blue")

            # 오른쪽 컬럼: content_blocks 또는 bullets
            if right_col.get("content_blocks"):
                # 원본 스펙 변형 방지를 위해 사본에 위치를 지정
                right_blocks = []
                for block in right_col["content_blocks"]:
                    right_block = dict(block)
                    right_block["position"] = "right"
                    right_blocks.append(right_block)
                self._render_content_blocks(slide, right_blocks, start_top_pt=160)
            else:
                self._add_bullets(
                    slide,
                    right_col.get("bullets", []),
                    left_pt=480, top_pt=160, width_pt=400, height_pt=300
                )
        else:
            # columns가 없으면 bullets에서 [Left]/[Right] 파싱
            bullets = slide_data.get("bullets", [])
            left_bullets = []
            right_bullets = []

            for b in bullets:
                text = b if isinstance(b, str) else b.get("text", "")
                if text.startswith("[Left]"):
                    left_bullets.append(text.replace("[Left]", "").strip())
                elif text.startswith("[Right]"):
                    right_bullets.append(text.replace("[Right]", "").strip())
                else:
                    left_bullets.append(text)

            self._add_bullets(
                slide, left_bullets,
                left_pt=43, top_pt=130, width_pt=400, height_pt=330
            )
            self._add_bullets(
                slide, right_bullets,
                left_pt=480, top_pt=130, width_pt=400, height_pt=330
            )

        # 슬라이드 레벨 보조 블록(하단 서술 등)
        if slide_data.get("content_blocks"):
            extra_blocks = [b for b in slide_data.get("content_blocks", []) if isinstance(b, dict)]
            if extra_blocks:
                self._render_content_blocks(slide, extra_blocks, start_top_pt=392)

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_three_column(self, slide, slide_data: dict):
        """3컬럼 레이아웃 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        columns = slide_data.get("columns", [])

        col_width = 270
        col_positions = [43, 330, 617]

        for i, col in enumerate(columns[:3]):
            # 컬럼 헤딩
            tx_head = slide.shapes.add_textbox(
                Pt(col_positions[i]), Pt(125),
                Pt(col_width), Pt(30)
            )
            tf = tx_head.text_frame
            tf.text = col.get("heading", "")
            self._apply_text_style(tf, "governing", "primary_blue")

            # 컬럼 콘텐츠 (content_blocks 우선, 없으면 bullets)
            if col.get("content_blocks"):
                current_top = 160
                for block in col["content_blocks"]:
                    block_type = block.get("type", "bullets")

                    if block_type == "bullets":
                        bullets = block.get("bullets", [])
                        block_height = self._estimate_bullet_block_height(bullets)
                        self._add_bullets(
                            slide, bullets,
                            left_pt=col_positions[i], top_pt=current_top,
                            width_pt=col_width, height_pt=block_height
                        )
                        current_top += block_height + 14

                    elif block_type == "text":
                        self._render_text_block(
                            slide, block.get("text", ""),
                            left_pt=col_positions[i], top_pt=current_top, width_pt=col_width
                        )
                        current_top += 70

                    elif block_type == "callout":
                        self._render_callout(
                            slide, block.get("callout", {}),
                            left_pt=col_positions[i], top_pt=current_top, width_pt=col_width
                        )
                        current_top += 84

                    elif block_type == "kpi":
                        self._render_kpi(
                            slide, block.get("kpi", {}),
                            left_pt=col_positions[i], top_pt=current_top, width_pt=col_width
                        )
                        current_top += 90

                    elif block_type == "quote":
                        self._render_quote(
                            slide, block.get("quote", {}),
                            left_pt=col_positions[i], top_pt=current_top, width_pt=col_width
                        )
                        current_top += 120

                    elif block_type == "table":
                        self._render_table_placeholder(
                            slide, block.get("table", {}),
                            left_pt=col_positions[i], top_pt=current_top, width_pt=col_width
                        )
                        current_top += 150

                    elif block_type == "chart":
                        self._render_chart_placeholder(
                            slide, block.get("chart", {}),
                            left_pt=col_positions[i], top_pt=current_top, width_pt=col_width
                        )
                        current_top += 220

                    elif block_type == "image":
                        self._render_image_placeholder(
                            slide, block.get("image", {}),
                            left_pt=col_positions[i], top_pt=current_top, width_pt=col_width
                        )
                        current_top += 220
            else:
                self._add_bullets(
                    slide,
                    col.get("bullets", []),
                    left_pt=col_positions[i], top_pt=160,
                    width_pt=col_width, height_pt=300
                )

        # 슬라이드 레벨 보조 블록(하단 서술 등)
        if slide_data.get("content_blocks"):
            extra_blocks = [b for b in slide_data.get("content_blocks", []) if isinstance(b, dict)]
            if extra_blocks:
                self._render_content_blocks(slide, extra_blocks, start_top_pt=392)

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_comparison(self, slide, slide_data: dict):
        """비교 레이아웃 렌더링 (As-Is / To-Be 등)"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        columns = slide_data.get("columns", [])

        if len(columns) >= 2:
            # 왼쪽 (As-Is)
            left_col = columns[0]
            # 배경 박스
            shape_left = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Pt(43), Pt(125), Pt(400), Pt(340)
            )
            shape_left.fill.solid()
            shape_left.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
            shape_left.line.fill.background()

            # 헤딩
            tx_left_head = slide.shapes.add_textbox(Pt(53), Pt(135), Pt(380), Pt(30))
            tf = tx_left_head.text_frame
            tf.text = left_col.get("heading", "As-Is")
            self._apply_text_style(tf, "governing", "dark_blue")

            # 왼쪽 컬럼: content_blocks + bullets
            if left_col.get("content_blocks"):
                left_blocks = [dict(block, position="left") for block in left_col.get("content_blocks", [])]
                if left_col.get("bullets") and not any(block.get("type") == "bullets" for block in left_blocks):
                    left_blocks.append({"type": "bullets", "position": "left", "bullets": left_col.get("bullets", [])})
                self._render_content_blocks(slide, left_blocks, start_top_pt=175)
            else:
                self._add_bullets(
                    slide,
                    left_col.get("bullets", []),
                    left_pt=53, top_pt=175, width_pt=380, height_pt=280
                )

            # 오른쪽 (To-Be)
            right_col = columns[1]
            shape_right = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Pt(480), Pt(125), Pt(400), Pt(340)
            )
            shape_right.fill.solid()
            shape_right.fill.fore_color.rgb = hex_to_rgb(self._get_color("primary_blue"))
            shape_right.line.fill.background()

            # 헤딩
            tx_right_head = slide.shapes.add_textbox(Pt(490), Pt(135), Pt(380), Pt(30))
            tf = tx_right_head.text_frame
            tf.text = right_col.get("heading", "To-Be")
            self._apply_text_style(tf, "governing", "background")

            # 오른쪽 컬럼: content_blocks + bullets
            if right_col.get("content_blocks"):
                right_blocks = [dict(block, position="right") for block in right_col.get("content_blocks", [])]
                if right_col.get("bullets") and not any(block.get("type") == "bullets" for block in right_blocks):
                    right_blocks.append({"type": "bullets", "position": "right", "bullets": right_col.get("bullets", [])})
                self._render_content_blocks(slide, right_blocks, start_top_pt=175)
            else:
                self._add_bullets(
                    slide,
                    right_col.get("bullets", []),
                    left_pt=490, top_pt=175, width_pt=380, height_pt=280,
                    color_key="background"
                )

            # 슬라이드 레벨 보조 블록(하단 서술 등)
            if slide_data.get("content_blocks"):
                extra_blocks = [b for b in slide_data.get("content_blocks", []) if isinstance(b, dict)]
                if extra_blocks:
                    self._render_content_blocks(slide, extra_blocks, start_top_pt=392)
        else:
            # columns 데이터가 없는 comparison도 content_blocks/bullets를 상단부터 렌더링
            blocks = [b for b in slide_data.get("content_blocks", []) if isinstance(b, dict)]
            bullets = slide_data.get("bullets", [])
            if bullets and not any(block.get("type") == "bullets" for block in blocks):
                blocks.insert(0, {"type": "bullets", "bullets": bullets})
            if blocks:
                self._render_content_blocks(slide, blocks, start_top_pt=160)

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_chart_focus(self, slide, slide_data: dict):
        """차트 중심 슬라이드 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        intent = slide_data.get("layout_intent", {})
        visual_position = str(intent.get("visual_position", "left")).lower()
        emphasis = str(intent.get("emphasis", "content")).lower()

        visual = self._extract_chart_visual(slide_data)

        # visual_position / emphasis를 반영한 배치
        if visual_position == "right":
            chart_left, chart_top, chart_width = 320, 130, 600
            bullet_left, bullet_top, bullet_width, bullet_height = 43, 130, 250, 330
        elif visual_position == "center" or emphasis == "visual":
            chart_left, chart_top, chart_width = 43, 130, 860
            bullet_left, bullet_top, bullet_width, bullet_height = 43, 360, 860, 120
        else:
            chart_left, chart_top, chart_width = 43, 130, 600
            bullet_left, bullet_top, bullet_width, bullet_height = 670, 130, 250, 330

        if visual:
            self._render_chart_placeholder(slide, visual, chart_left, chart_top, chart_width)

        bullets = slide_data.get("bullets", [])
        aux_blocks = [
            block for block in slide_data.get("content_blocks", [])
            if isinstance(block, dict) and block.get("type") != "chart"
        ]

        side_blocks = []
        if bullets:
            side_blocks.append({"type": "bullets", "bullets": bullets})
        side_blocks.extend(aux_blocks)

        if side_blocks:
            # 시각 요소 옆 설명 패널
            if bullet_width <= 280:
                panel = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Pt(bullet_left), Pt(bullet_top), Pt(bullet_width), Pt(bullet_height)
                )
                panel.fill.solid()
                panel.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
                panel.line.color.rgb = hex_to_rgb(self._get_color("divider_gray"))

            stacked_blocks = []
            for block in side_blocks:
                block_copy = dict(block)
                block_copy["position"] = "main"
                block_copy["left_pt"] = bullet_left + 8
                block_copy["width_pt"] = max(140, bullet_width - 16)
                stacked_blocks.append(block_copy)

            self._render_content_blocks(
                slide,
                stacked_blocks,
                start_top_pt=bullet_top + 8
            )

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_image_focus(self, slide, slide_data: dict):
        """이미지 중심 슬라이드 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        intent = slide_data.get("layout_intent", {})
        visual_position = str(intent.get("visual_position", "left")).lower()
        emphasis = str(intent.get("emphasis", "content")).lower()

        visual = self._extract_image_visual(slide_data)

        if visual_position == "right":
            image_left, image_top, image_width = 320, 130, 600
            bullet_left, bullet_top, bullet_width, bullet_height = 43, 130, 250, 330
        elif visual_position == "center" or emphasis == "visual":
            image_left, image_top, image_width = 43, 130, 860
            bullet_left, bullet_top, bullet_width, bullet_height = 43, 360, 860, 120
        else:
            image_left, image_top, image_width = 43, 130, 600
            bullet_left, bullet_top, bullet_width, bullet_height = 670, 130, 250, 330

        if visual:
            self._render_image_placeholder(slide, visual, image_left, image_top, image_width)

        bullets = slide_data.get("bullets", [])
        aux_blocks = [
            block for block in slide_data.get("content_blocks", [])
            if isinstance(block, dict) and block.get("type") != "image"
        ]

        side_blocks = []
        if bullets:
            side_blocks.append({"type": "bullets", "bullets": bullets})
        side_blocks.extend(aux_blocks)

        if side_blocks:
            if bullet_width <= 280:
                panel = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Pt(bullet_left), Pt(bullet_top), Pt(bullet_width), Pt(bullet_height)
                )
                panel.fill.solid()
                panel.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
                panel.line.color.rgb = hex_to_rgb(self._get_color("divider_gray"))

            stacked_blocks = []
            for block in side_blocks:
                block_copy = dict(block)
                block_copy["position"] = "main"
                block_copy["left_pt"] = bullet_left + 8
                block_copy["width_pt"] = max(140, bullet_width - 16)
                stacked_blocks.append(block_copy)

            self._render_content_blocks(
                slide,
                stacked_blocks,
                start_top_pt=bullet_top + 8
            )

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_quote_layout(self, slide, slide_data: dict):
        """인용문 레이아웃 렌더링"""
        title = slide_data.get("title", "")
        governing = slide_data.get("governing_message", "")

        # content_blocks에서 quote 찾기
        content_blocks = slide_data.get("content_blocks", [])
        quote_block = None
        for block in content_blocks:
            if block.get("type") == "quote":
                quote_block = block.get("quote", {})
                break

        # 중앙 제목
        if title:
            tx = slide.shapes.add_textbox(Pt(43), Pt(150), Pt(860), Pt(50))
            tf = tx.text_frame
            tf.text = title
            tf.word_wrap = True
            self._apply_text_style(tf, "title", "primary_blue", PP_ALIGN.CENTER)

        # 인용문
        if quote_block:
            self._render_quote(slide, quote_block, 100, 230, 750)
        elif governing:
            # governing을 인용문처럼 표시
            quote_def = {"text": governing}
            self._render_quote(slide, quote_def, 100, 230, 750)

        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_timeline(self, slide, slide_data: dict):
        """타임라인 슬라이드 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        # 타임라인 화살표 기본 도형
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Pt(43), Pt(250), Pt(860), Pt(40)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = hex_to_rgb(self._get_color("primary_blue"))
        arrow.line.fill.background()

        # 불릿을 타임라인 단계로 표시
        bullets = slide_data.get("bullets", [])
        if bullets:
            step_width = 860 // len(bullets)
            for i, bullet in enumerate(bullets):
                text = bullet if isinstance(bullet, str) else bullet.get("text", "")

                # 단계 박스
                tx = slide.shapes.add_textbox(
                    Pt(43 + i * step_width), Pt(300),
                    Pt(step_width - 10), Pt(80)
                )
                tf = tx.text_frame
                tf.text = f"Phase {i + 1}\n{text}"
                tf.word_wrap = True
                self._apply_text_style(tf, "body", "text_dark", PP_ALIGN.CENTER)

        if slide_data.get("content_blocks"):
            extra_blocks = [b for b in slide_data.get("content_blocks", []) if isinstance(b, dict)]
            if extra_blocks:
                self._render_content_blocks(slide, extra_blocks, start_top_pt=392)

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_process_flow(self, slide, slide_data: dict):
        """프로세스 플로우 슬라이드 렌더링"""
        self._set_title(slide, slide_data.get("title", ""))
        self._add_governing_message(slide, slide_data.get("governing_message", ""))

        steps = slide_data.get("bullets", [])
        if not steps:
            self._add_notes(slide, slide_data.get("notes", ""))
            return

        max_steps = min(len(steps), 5)
        steps = steps[:max_steps]
        total_width = 860
        gap = 14
        box_width = int((total_width - (gap * (max_steps - 1))) / max_steps)
        base_left = 43
        box_top = 200
        box_height = 130

        for i, step in enumerate(steps):
            step_text = step if isinstance(step, str) else step.get("text", "")
            box_left = base_left + (i * (box_width + gap))

            # 단계 박스
            step_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Pt(box_left), Pt(box_top), Pt(box_width), Pt(box_height)
            )
            step_shape.fill.solid()
            step_shape.fill.fore_color.rgb = hex_to_rgb(self._get_color("light_blue"))
            step_shape.line.color.rgb = hex_to_rgb(self._get_color("primary_blue"))

            tf = step_shape.text_frame
            tf.word_wrap = True
            tf.text = f"Step {i + 1}\n{step_text}"
            self._apply_text_style(tf, "body", "text_dark", PP_ALIGN.CENTER)

            # 박스 사이 화살표
            if i < max_steps - 1:
                arrow_left = box_left + box_width + 2
                arrow_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Pt(arrow_left), Pt(box_top + 45), Pt(gap - 4), Pt(30)
                )
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = hex_to_rgb(self._get_color("primary_blue"))
                arrow_shape.line.fill.background()

        if slide_data.get("content_blocks"):
            extra_blocks = [b for b in slide_data.get("content_blocks", []) if isinstance(b, dict)]
            if extra_blocks:
                self._render_content_blocks(slide, extra_blocks, start_top_pt=370)

        self._add_footnotes(slide, slide_data.get("footnotes", []))
        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_thank_you(self, slide, slide_data: dict):
        """마무리 슬라이드 렌더링"""
        title = slide_data.get("title", "Thank You")
        governing = slide_data.get("governing_message", "")

        # 중앙 큰 제목
        tx = slide.shapes.add_textbox(Pt(43), Pt(200), Pt(860), Pt(80))
        tf = tx.text_frame
        tf.text = title
        tf.word_wrap = True
        self._apply_text_style(tf, "title", "primary_blue", PP_ALIGN.CENTER)

        # 서브 메시지
        if governing:
            tx2 = slide.shapes.add_textbox(Pt(43), Pt(290), Pt(860), Pt(50))
            tf2 = tx2.text_frame
            tf2.text = governing
            tf2.word_wrap = True
            self._apply_text_style(tf2, "governing", "text_muted", PP_ALIGN.CENTER)

        self._add_notes(slide, slide_data.get("notes", ""))

    def _render_appendix(self, slide, slide_data: dict):
        """부록 슬라이드 렌더링"""
        self._render_content(slide, slide_data)

    def _render_default(self, slide, slide_data: dict):
        """기본 렌더링 (알 수 없는 레이아웃)"""
        self._render_content(slide, slide_data)

    # =========================================================================
    # 메인 렌더링 로직
    # =========================================================================

    def render(
        self,
        spec: dict,
        template_path: Path,
        output_path: Path,
        spec_base_dir: Optional[Path] = None
    ):
        """덱 스펙을 PPTX로 렌더링"""
        self.asset_base_dir = spec_base_dir

        # 프레젠테이션 로드/생성
        if template_path.exists():
            self.prs = Presentation(str(template_path))
            # 기존 슬라이드 제거 (템플릿의 슬라이드 마스터/레이아웃만 유지)
            # 슬라이드는 역순으로 제거해야 인덱스 문제 방지
            while len(self.prs.slides) > 0:
                rId = self.prs.slides._sldIdLst[0].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[0]
        else:
            self.prs = Presentation()
            # 템플릿 미사용 시에도 16:9 와이드 규격으로 통일
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            print(f"Warning: 템플릿 없음, 빈 프레젠테이션 사용: {template_path}")

        # 레이아웃별 렌더러 매핑
        layout_renderers = {
            "cover": self._render_cover,
            "exec_summary": self._render_exec_summary,
            "section_divider": self._render_section_divider,
            "content": self._render_content,
            "two_column": self._render_two_column,
            "three_column": self._render_three_column,
            "comparison": self._render_comparison,
            "timeline": self._render_timeline,
            "process_flow": self._render_process_flow,
            "chart_focus": self._render_chart_focus,
            "image_focus": self._render_image_focus,
            "quote": self._render_quote_layout,
            "appendix": self._render_appendix,
            "thank_you": self._render_thank_you,
        }

        # 슬라이드 렌더링
        for slide_data in spec.get("slides", []):
            layout_name = slide_data.get("layout", "content")

            # 슬라이드 레이아웃 선택
            slide_layout = self._pick_layout(layout_name)
            slide = self.prs.slides.add_slide(slide_layout)

            # 레이아웃별 렌더러 호출
            renderer = layout_renderers.get(layout_name, self._render_default)
            renderer(slide, slide_data)
            self._add_slide_chrome(slide, layout_name)

        # 저장
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))


def render(
    spec_path: Path,
    template_pptx: Path,
    output_pptx: Path,
    tokens_path: Path,
    layouts_path: Path
):
    """렌더링 실행 함수 (외부 호출용)"""
    spec = load_yaml(spec_path)
    tokens = load_yaml(tokens_path)
    layouts = load_yaml(layouts_path)

    renderer = DeckRenderer(tokens, layouts)
    renderer.render(spec, template_pptx, output_pptx, spec_base_dir=spec_path.parent)


def main():
    if len(sys.argv) != 5:
        print("Usage: python scripts/render_ppt.py <deck_spec.yaml> <template.pptx> <output.pptx> <templates_dir>")
        print("Example: python scripts/render_ppt.py clients/acme-demo/deck_spec.yaml templates/company/base-template.pptx clients/acme-demo/outputs/acme-demo.pptx templates/company")
        sys.exit(1)

    spec_path = Path(sys.argv[1]).resolve()
    template_pptx = Path(sys.argv[2]).resolve()
    output_pptx = Path(sys.argv[3]).resolve()
    tokens_dir = Path(sys.argv[4]).resolve()

    tokens_path = tokens_dir / "tokens.yaml"
    layouts_path = tokens_dir / "layouts.yaml"

    # 파일 존재 확인
    if not spec_path.exists():
        print(f"Error: deck_spec not found: {spec_path}")
        sys.exit(1)
    if not tokens_path.exists():
        print(f"Error: tokens.yaml not found: {tokens_path}")
        sys.exit(1)
    if not layouts_path.exists():
        print(f"Error: layouts.yaml not found: {layouts_path}")
        sys.exit(1)

    render(spec_path, template_pptx, output_pptx, tokens_path, layouts_path)
    print(f"✓ Rendered PPTX: {output_pptx}")


if __name__ == "__main__":
    main()
