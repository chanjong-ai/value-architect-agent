#!/usr/bin/env python3
"""
polish_ppt.py - 렌더링된 PPTX 미세 편집기

기능:
1. 폰트 일관성 정리 (tokens.yaml 기준)
2. 텍스트 공백/탭 정리
3. 본문 문단 줄간격 정리
4. 편집 로그(JSON) 생성
"""

import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Optional

import yaml

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx 패키지가 필요합니다: pip install python-pptx")
    sys.exit(1)


def load_yaml(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        return yaml.safe_load(f) or {}


def get_font_name(tokens: dict, key: str, fallback: str) -> str:
    font_cfg = tokens.get("fonts", {}).get(key, {})
    if isinstance(font_cfg, dict):
        return font_cfg.get("name") or font_cfg.get("family") or fallback
    if isinstance(font_cfg, str):
        return font_cfg
    return fallback


def normalize_run_text(text: str) -> str:
    # 탭 제거, 연속 공백 축소
    text = text.replace("\t", " ")
    text = re.sub(r" {2,}", " ", text)
    return text


def polish_ppt(
    input_pptx: Path,
    output_pptx: Path,
    tokens_path: Optional[Path] = None,
    report_path: Optional[Path] = None
) -> dict:
    prs = Presentation(str(input_pptx))
    tokens = load_yaml(tokens_path) if tokens_path and tokens_path.exists() else {}

    render_options = tokens.get("render_options", {}) if isinstance(tokens.get("render_options", {}), dict) else {}
    preserve_template_fonts = bool(render_options.get("preserve_template_fonts_in_polish", True))
    preserve_template_line_spacing = bool(render_options.get("preserve_template_line_spacing_in_polish", False))
    force_regular_weight = bool(render_options.get("force_regular_weight_in_polish", True))

    title_font = get_font_name(tokens, "title", "Noto Sans KR")
    governing_font = get_font_name(tokens, "governing", "Noto Sans KR")
    body_font = get_font_name(tokens, "body", "Noto Sans KR")
    footnote_font = get_font_name(tokens, "footnote", "Noto Sans KR")

    stats = {
        "slides": len(prs.slides),
        "font_updates": 0,
        "text_normalizations": 0,
        "line_spacing_updates": 0,
        "weight_normalizations": 0,
    }
    changes = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            for para_idx, para in enumerate(tf.paragraphs, start=1):
                # 문단 위치/크기 기반으로 폰트 타입 추정
                target_font = body_font
                if hasattr(shape, "top"):
                    if shape.top < prs.slide_height * 0.12:
                        target_font = title_font
                    elif shape.top < prs.slide_height * 0.22:
                        target_font = governing_font
                    elif shape.top > prs.slide_height * 0.86:
                        target_font = footnote_font

                # 본문 계열 문단 줄간격 정리
                if (not preserve_template_line_spacing) and target_font in (body_font, governing_font):
                    para.line_spacing = 1.15
                    stats["line_spacing_updates"] += 1

                for run in para.runs:
                    before_font = run.font.name
                    before_text = run.text

                    # 폰트 통일
                    if (not preserve_template_fonts) and before_font != target_font:
                        run.font.name = target_font
                        stats["font_updates"] += 1

                    # 사용자 요청: 전체 일반체 강제
                    if force_regular_weight and run.font.bold is not False:
                        run.font.bold = False
                        stats["weight_normalizations"] += 1

                    # 공백 정리
                    after_text = normalize_run_text(before_text)
                    if after_text != before_text:
                        run.text = after_text
                        stats["text_normalizations"] += 1

                if para.runs:
                    changes.append({
                        "slide": slide_idx,
                        "paragraph": para_idx,
                        "target_font": target_font,
                    })

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))

    result = {
        "timestamp": datetime.now().isoformat(timespec="seconds"),
        "source_pptx": str(input_pptx),
        "output_pptx": str(output_pptx),
        "stats": stats,
        "changes_preview": changes[:50],
    }

    if report_path:
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")

    return result


def main():
    parser = argparse.ArgumentParser(description="PPTX 미세 편집기")
    parser.add_argument("input_pptx", help="입력 PPTX 파일")
    parser.add_argument("--output", "-o", help="출력 PPTX 파일")
    parser.add_argument("--tokens", help="tokens.yaml 파일 경로")
    parser.add_argument("--report", help="편집 로그(JSON) 파일 경로")
    args = parser.parse_args()

    input_pptx = Path(args.input_pptx).resolve()
    if not input_pptx.exists():
        print(f"Error: input file not found: {input_pptx}")
        return 1

    if args.output:
        output_pptx = Path(args.output).resolve()
    else:
        output_pptx = input_pptx.with_name(f"{input_pptx.stem}_polished.pptx")

    tokens_path = Path(args.tokens).resolve() if args.tokens else None
    report_path = Path(args.report).resolve() if args.report else output_pptx.with_suffix(".polish.json")

    result = polish_ppt(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        tokens_path=tokens_path,
        report_path=report_path
    )

    print(f"✓ Polished PPTX: {result['output_pptx']}")
    print(
        "  - 폰트 변경: {font_updates}, 텍스트 정리: {text_normalizations}, 줄간격 조정: {line_spacing_updates}, weight 정리: {weight_normalizations}".format(
            **result["stats"]
        )
    )
    print(f"  - 로그: {report_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
