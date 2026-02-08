#!/usr/bin/env python3
"""
PPT QA (Quality Assurance) 자동 검사기 v2.1

개선사항:
1. tokens.yaml 파싱 규칙 수정 (name/size_pt 구조 + 레거시 family/size 호환)
2. 불릿 카운트에서 제목/거버닝 메시지 텍스트 박스 제외
3. columns[].bullets도 불릿 수/길이 규칙에 포함
4. evidence 검사: sources.md# 앵커 포맷 검증 및 존재 검사
5. global_constraints + slide_constraints 적용 (로컬 오버라이드 반영)
6. Spec alignment 탐지 로직 고도화 (폰트 크기/위치 기준)
7. 불릿 길이 기준 단일 정책 (constants.py 사용)

사용법:
    python qa_ppt.py <pptx_path> [--spec <spec_path>] [--tokens <tokens_path>] [--sources <sources_path>] [--output <report_path>]
"""

import argparse
import json
import re
import sys
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional, List, Tuple
from enum import Enum

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("python-pptx 패키지가 필요합니다: pip install python-pptx")
    sys.exit(1)

try:
    import yaml
except ImportError:
    print("PyYAML 패키지가 필요합니다: pip install pyyaml")
    sys.exit(1)

# 상수 모듈 임포트
try:
    from constants import (
        BULLET_MAX_CHARS, BULLET_MAX_COUNT, BULLET_MIN_COUNT,
        BULLET_CHARS_PER_LINE, BULLET_MAX_LINES,
        BULLET_RECOMMENDED_MIN_CHARS, BULLET_RECOMMENDED_MAX_CHARS,
        BULLET_BLOCK_MIN_ITEMS, BULLET_BLOCK_MAX_ITEMS,
        ACTION_LIST_MIN_ITEMS, ACTION_LIST_MAX_ITEMS,
        TITLE_FONT_SIZE_PT, GOVERNING_FONT_SIZE_PT, BODY_FONT_SIZE_PT,
        FONT_SIZE_TOLERANCE_PT, ALLOWED_FONTS,
        DENSITY_MAX_CHARS, DENSITY_MIN_CHARS, DENSITY_MIN_PARAGRAPHS,
        NO_BULLET_LAYOUTS, COLUMN_LAYOUTS, EVIDENCE_ANCHOR_PATTERN,
        get_max_bullets, get_max_chars_per_bullet, get_forbidden_words, get_bullet_bounds,
        get_column_bullet_limit, normalize_layout_name, LAYOUT_REQUIRED_BLOCKS
    )
except ImportError:
    # 폴백 상수
    BULLET_MAX_CHARS = 180
    BULLET_MAX_COUNT = 9
    BULLET_MIN_COUNT = 3
    BULLET_CHARS_PER_LINE = 38
    BULLET_MAX_LINES = 4
    BULLET_RECOMMENDED_MIN_CHARS = 18
    BULLET_RECOMMENDED_MAX_CHARS = 110
    BULLET_BLOCK_MIN_ITEMS = 3
    BULLET_BLOCK_MAX_ITEMS = 5
    ACTION_LIST_MIN_ITEMS = 2
    ACTION_LIST_MAX_ITEMS = 3
    TITLE_FONT_SIZE_PT = 24
    GOVERNING_FONT_SIZE_PT = 16
    BODY_FONT_SIZE_PT = 12
    FONT_SIZE_TOLERANCE_PT = 2
    ALLOWED_FONTS = ["Noto Sans KR", "NotoSansKR"]
    DENSITY_MAX_CHARS = 1200
    DENSITY_MIN_CHARS = 50
    DENSITY_MIN_PARAGRAPHS = 3
    NO_BULLET_LAYOUTS = ["cover", "section_divider", "thank_you", "quote"]
    COLUMN_LAYOUTS = ["two_column", "three_column", "comparison"]
    EVIDENCE_ANCHOR_PATTERN = r"^sources\.md#[\w-]+$"
    LAYOUT_REQUIRED_BLOCKS = {
        "exec_summary": ["bullets", "action_list"],
        "two_column": ["bullets", "action_list"],
        "chart_insight": ["chart", "bullets", "action_list"],
        "competitor_2x2": ["matrix_2x2", "bullets", "action_list"],
        "strategy_cards": ["kpi_cards", "action_list"],
        "timeline": ["timeline_steps", "action_list"],
        "kpi_cards": ["kpi_cards", "action_list"],
    }

    def get_max_bullets(gc=None, sc=None):
        if sc and "max_bullets" in sc:
            return sc["max_bullets"]
        if gc and "default_max_bullets" in gc:
            return gc["default_max_bullets"]
        return BULLET_MAX_COUNT

    def get_max_chars_per_bullet(gc=None, sc=None):
        if sc and "max_chars_per_bullet" in sc:
            return sc["max_chars_per_bullet"]
        if gc and "default_max_chars_per_bullet" in gc:
            return gc["default_max_chars_per_bullet"]
        return BULLET_MAX_CHARS

    def get_forbidden_words(gc=None, sc=None):
        words = []
        if gc and "forbidden_words" in gc:
            words.extend(gc["forbidden_words"])
        if sc and "forbidden_words" in sc:
            words.extend(sc["forbidden_words"])
        return list(set(words))

    def get_bullet_bounds(layout, gc=None, sc=None):
        layout = normalize_layout_name(layout)
        max_bullets = get_max_bullets(gc, sc)
        min_bullets = BULLET_MIN_COUNT
        if layout in NO_BULLET_LAYOUTS:
            return 0, 0
        if layout in ("chart_focus", "image_focus", "chart_insight", "competitor_2x2", "kpi_cards"):
            return 0, min(max_bullets, 8)
        return min_bullets, max_bullets

    def get_column_bullet_limit(max_bullets):
        if max_bullets <= 0:
            return 0
        return max(3, min(8, max_bullets))

    def normalize_layout_name(layout):
        key = str(layout or "").strip().lower()
        return {"chart_focus": "chart_insight", "strategy_options": "strategy_cards"}.get(key, key)

try:
    from block_utils import normalize_slide_blocks, block_types_in_slide, iter_bullet_texts
except ImportError:
    normalize_slide_blocks = None
    block_types_in_slide = None
    iter_bullet_texts = None


class Severity(Enum):
    """QA 이슈 심각도"""
    ERROR = "error"          # 반드시 수정 필요
    WARNING = "warning"      # 권장 수정
    INFO = "info"            # 참고 사항


@dataclass
class QAIssue:
    """QA 이슈"""
    slide_index: int
    severity: Severity
    category: str
    message: str
    details: dict = field(default_factory=dict)
    auto_fixable: bool = False


@dataclass
class QAReport:
    """QA 보고서"""
    pptx_path: str
    total_slides: int
    issues: list[QAIssue] = field(default_factory=list)
    summary: dict = field(default_factory=dict)

    def add_issue(self, issue: QAIssue):
        self.issues.append(issue)

    @property
    def error_count(self) -> int:
        return sum(1 for i in self.issues if i.severity == Severity.ERROR)

    @property
    def warning_count(self) -> int:
        return sum(1 for i in self.issues if i.severity == Severity.WARNING)

    @property
    def info_count(self) -> int:
        return sum(1 for i in self.issues if i.severity == Severity.INFO)

    @property
    def passed(self) -> bool:
        return self.error_count == 0

    def to_dict(self) -> dict:
        return {
            "pptx_path": self.pptx_path,
            "total_slides": self.total_slides,
            "passed": self.passed,
            "summary": {
                "errors": self.error_count,
                "warnings": self.warning_count,
                "info": self.info_count
            },
            "issues": [
                {
                    "slide": i.slide_index,
                    "severity": i.severity.value,
                    "category": i.category,
                    "message": i.message,
                    "details": i.details,
                    "auto_fixable": i.auto_fixable
                }
                for i in self.issues
            ]
        }

    def to_markdown(self) -> str:
        """마크다운 형식 보고서"""
        lines = [
            "# PPT QA 보고서",
            "",
            f"**파일**: `{self.pptx_path}`",
            f"**슬라이드 수**: {self.total_slides}",
            f"**결과**: {'✅ 통과' if self.passed else '❌ 실패'}",
            "",
            "## 요약",
            f"- 🔴 오류: {self.error_count}",
            f"- 🟡 경고: {self.warning_count}",
            f"- 🔵 참고: {self.info_count}",
            ""
        ]

        if self.issues:
            lines.append("## 상세 이슈")
            lines.append("")

            # 슬라이드별로 그룹핑
            by_slide: dict[int, list[QAIssue]] = {}
            for issue in self.issues:
                if issue.slide_index not in by_slide:
                    by_slide[issue.slide_index] = []
                by_slide[issue.slide_index].append(issue)

            for slide_idx in sorted(by_slide.keys()):
                if slide_idx == 0:
                    lines.append("### 전역 이슈")
                else:
                    lines.append(f"### 슬라이드 {slide_idx}")
                for issue in by_slide[slide_idx]:
                    icon = {"error": "🔴", "warning": "🟡", "info": "🔵"}[issue.severity.value]
                    fix_tag = " [자동수정가능]" if issue.auto_fixable else ""
                    lines.append(f"- {icon} **[{issue.category}]** {issue.message}{fix_tag}")
                    if issue.details:
                        for k, v in issue.details.items():
                            lines.append(f"  - {k}: {v}")
                lines.append("")
        else:
            lines.append("✅ 모든 검사 통과!")

        return "\n".join(lines)


class PPTQAChecker:
    """PPT QA 검사기 v2.1"""

    def __init__(self, pptx_path: str, spec_path: Optional[str] = None,
                 tokens_path: Optional[str] = None, sources_path: Optional[str] = None,
                 layouts_path: Optional[str] = None):
        self.pptx_path = Path(pptx_path)
        self.spec_path = Path(spec_path) if spec_path else None
        self.tokens_path = Path(tokens_path) if tokens_path else None
        self.sources_path = Path(sources_path) if sources_path else None
        self.layouts_path = Path(layouts_path) if layouts_path else None

        self.prs = Presentation(str(self.pptx_path))
        self.spec = self._load_spec() if self.spec_path else None
        self.tokens = self._load_tokens() if self.tokens_path else None
        self.layouts = self._load_layouts() if self.layouts_path else None
        self.sources_anchors = self._parse_sources_anchors() if self.sources_path else set()

        # 전역 제약조건
        self.global_constraints = self.spec.get("global_constraints", {}) if self.spec else {}

        # 제약조건 빌드
        self.constraints = self._build_constraints()

        self.report = QAReport(
            pptx_path=str(self.pptx_path),
            total_slides=len(self.prs.slides)
        )

    def _load_spec(self) -> Optional[dict]:
        """deck_spec.yaml 로드"""
        if not self.spec_path or not self.spec_path.exists():
            return None
        with open(self.spec_path, 'r', encoding='utf-8') as f:
            return yaml.safe_load(f)

    def _load_tokens(self) -> Optional[dict]:
        """tokens.yaml 로드"""
        if not self.tokens_path or not self.tokens_path.exists():
            return None
        with open(self.tokens_path, 'r', encoding='utf-8') as f:
            return yaml.safe_load(f)

    def _parse_sources_anchors(self) -> set:
        """sources.md에서 앵커 파싱"""
        anchors = set()
        if not self.sources_path or not self.sources_path.exists():
            return anchors

        with open(self.sources_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # 마크다운 헤딩에서 앵커 추출 (# Heading -> heading)
        # GitHub 스타일 앵커: 소문자, 공백은 하이픈
        heading_pattern = r'^#+\s+(.+)$'
        for match in re.finditer(heading_pattern, content, re.MULTILINE):
            heading = match.group(1).strip()
            # 앵커 변환: 소문자, 공백->하이픈, 특수문자 제거
            anchor = re.sub(r'[^\w\s-]', '', heading.lower())
            anchor = re.sub(r'[\s]+', '-', anchor)
            anchors.add(f"sources.md#{anchor}")

        return anchors

    def _load_layouts(self) -> Optional[dict]:
        """layouts.yaml 로드"""
        if not self.layouts_path or not self.layouts_path.exists():
            return None
        with open(self.layouts_path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f) or {}

    def _get_font_from_tokens(self, font_key: str) -> Tuple[str, int]:
        """
        tokens.yaml에서 폰트 정보 추출
        개선: name/size_pt 구조 및 레거시 family/size 호환
        """
        if not self.tokens:
            return ("Noto Sans KR", BODY_FONT_SIZE_PT)

        fonts = self.tokens.get("fonts", {})
        font_info = fonts.get(font_key, {})

        if isinstance(font_info, dict):
            # 신규 구조: name/size_pt
            name = font_info.get("name") or font_info.get("family", "Noto Sans KR")
            size = font_info.get("size_pt") or font_info.get("size", BODY_FONT_SIZE_PT)
            return (name, size)
        elif isinstance(font_info, str):
            # 단순 문자열인 경우 (레거시)
            return (font_info, BODY_FONT_SIZE_PT)

        return ("Noto Sans KR", BODY_FONT_SIZE_PT)

    def _build_constraints(self) -> dict:
        """제약조건 빌드 (tokens + global_constraints 병합)"""
        constraints = {
            "max_bullets_per_slide": BULLET_MAX_COUNT,
            "max_chars_per_bullet": BULLET_MAX_CHARS,
            "min_bullets_per_content_slide": BULLET_MIN_COUNT,
            "max_title_chars": 100,
            "max_governing_chars": 200,
            "title_font_size_pt": TITLE_FONT_SIZE_PT,
            "governing_font_size_pt": GOVERNING_FONT_SIZE_PT,
            "body_font_size_pt": BODY_FONT_SIZE_PT,
            "font_size_tolerance_pt": FONT_SIZE_TOLERANCE_PT,
            "allowed_fonts": ALLOWED_FONTS.copy(),
            "forbidden_words": [],
        }

        # tokens.yaml에서 폰트 정보 로드 (개선된 파싱)
        if self.tokens:
            fonts = self.tokens.get("fonts", {})
            if fonts:
                # 허용 폰트 목록 구축
                allowed = []
                for key in ["title", "governing", "body", "footnote"]:
                    if key in fonts:
                        name, _ = self._get_font_from_tokens(key)
                        if name:
                            allowed.append(name)
                if allowed:
                    # 기존 목록과 병합
                    constraints["allowed_fonts"] = list(set(allowed + constraints["allowed_fonts"]))

                # 폰트 사이즈 업데이트
                _, title_size = self._get_font_from_tokens("title")
                _, governing_size = self._get_font_from_tokens("governing")
                _, body_size = self._get_font_from_tokens("body")

                constraints["title_font_size_pt"] = title_size
                constraints["governing_font_size_pt"] = governing_size
                constraints["body_font_size_pt"] = body_size

            # density_rules에서 불릿 규칙 로드
            density = self.tokens.get("density_rules", {})
            if density:
                if "bullets_min" in density:
                    constraints["min_bullets_per_content_slide"] = density["bullets_min"]
                if "bullets_max" in density:
                    constraints["max_bullets_per_slide"] = density["bullets_max"]

        # global_constraints 적용
        if self.global_constraints:
            gc = self.global_constraints
            if "default_max_bullets" in gc:
                constraints["max_bullets_per_slide"] = gc["default_max_bullets"]
            if "default_max_chars_per_bullet" in gc:
                constraints["max_chars_per_bullet"] = gc["default_max_chars_per_bullet"]
            if "forbidden_words" in gc:
                constraints["forbidden_words"] = gc["forbidden_words"]

        return constraints

    def _get_slide_constraints(self, slide_idx: int) -> dict:
        """슬라이드별 제약조건 가져오기 (global + slide_constraints 병합)"""
        slide_constraints = {}
        if self.spec and "slides" in self.spec:
            slides = self.spec["slides"]
            if 0 <= slide_idx - 1 < len(slides):
                spec_slide = slides[slide_idx - 1]
                slide_constraints = spec_slide.get("slide_constraints", {})

        return slide_constraints

    def run_all_checks(self) -> QAReport:
        """모든 QA 검사 실행"""
        for slide_idx, slide in enumerate(self.prs.slides, start=1):
            self._check_slide(slide_idx, slide)

        # 전역 검사
        self._check_global()

        # Evidence 검사 (spec 기준)
        if self.spec:
            self._check_evidence()

        return self.report

    def _check_slide(self, slide_idx: int, slide):
        """개별 슬라이드 검사"""
        # 슬라이드별 제약조건 가져오기
        slide_constraints = self._get_slide_constraints(slide_idx)

        # 1. 불릿 검사 (제목/거버닝 메시지 제외)
        self._check_bullets(slide_idx, slide, slide_constraints)

        # 2. 폰트 검사
        self._check_fonts(slide_idx, slide)

        # 3. 콘텐츠 밀도 검사
        self._check_density(slide_idx, slide)

        # 4. 레이아웃 경계 검사
        self._check_layout_bounds(slide_idx, slide)

        # 5. 레이아웃 필수 블록 검사 (spec 기반)
        self._check_required_blocks(slide_idx)

        # 6. 블록 구조/밀도 규칙 검사 (spec 기반)
        self._check_block_density_rules(slide_idx)

        # 7. 텍스트 오버플로우 추정
        self._check_text_overflow(slide_idx, slide)

        # 8. 금지어 검사 (slide_constraints 포함)
        self._check_forbidden_words(slide_idx, slide, slide_constraints)

        # 9. Spec과의 일치 검사 (고도화된 로직)
        if self.spec:
            self._check_spec_alignment(slide_idx, slide)

    def _extract_spec_bullet_texts(self, slide_idx: int) -> List[str]:
        """Spec에서 슬라이드의 모든 불릿 텍스트를 추출"""
        if not self.spec or "slides" not in self.spec:
            return []

        slides = self.spec["slides"]
        if slide_idx - 1 >= len(slides):
            return []

        spec_slide = slides[slide_idx - 1]
        bullet_texts: List[str] = []

        def _append_bullets(items):
            for item in items:
                if isinstance(item, str):
                    text = item.strip()
                elif isinstance(item, dict):
                    text = str(item.get("text", "")).strip()
                else:
                    text = ""
                if text:
                    bullet_texts.append(text)

        # Top-level bullets
        _append_bullets(spec_slide.get("bullets", []))

        # Columns bullets + column content_blocks bullets
        for col in spec_slide.get("columns", []):
            _append_bullets(col.get("bullets", []))
            for block in col.get("content_blocks", []):
                if block.get("type") == "bullets":
                    _append_bullets(block.get("bullets", []))

        # Slide-level content_blocks bullets
        for block in spec_slide.get("content_blocks", []):
            if block.get("type") == "bullets":
                _append_bullets(block.get("bullets", []))

        # New blocks bullets/action_list
        normalized_blocks = normalize_slide_blocks(spec_slide) if normalize_slide_blocks else spec_slide.get("blocks", [])
        for block in normalized_blocks if isinstance(normalized_blocks, list) else []:
            if not isinstance(block, dict):
                continue
            b_type = str(block.get("type", "")).strip().lower()
            if b_type == "bullets":
                _append_bullets(block.get("items", []))

        return bullet_texts

    def _get_spec_slide(self, slide_idx: int) -> Optional[dict]:
        """슬라이드 번호에 대응하는 spec 슬라이드 반환"""
        if not self.spec or "slides" not in self.spec:
            return None
        slides = self.spec.get("slides", [])
        if 0 <= slide_idx - 1 < len(slides):
            return slides[slide_idx - 1]
        return None

    @staticmethod
    def _collect_column_bullet_texts_spec(column: dict) -> List[str]:
        """컬럼 내 bullets + bullets block 텍스트 수집"""
        texts: List[str] = []

        def _append(items):
            for item in items:
                if isinstance(item, str):
                    text = item.strip()
                elif isinstance(item, dict):
                    text = str(item.get("text", "")).strip()
                else:
                    text = ""
                if text:
                    texts.append(text)

        _append(column.get("bullets", []))
        for block in column.get("content_blocks", []):
            if block.get("type") == "bullets":
                _append(block.get("bullets", []))

        return texts

    @staticmethod
    def _column_has_non_bullet_content_spec(column: dict) -> bool:
        """컬럼에 bullets 외 콘텐츠가 있는지 검사"""
        for block in column.get("content_blocks", []):
            if str(block.get("type", "")).strip().lower() != "bullets":
                return True
        return False

    def _slide_has_non_bullet_content_spec(self, spec_slide: dict) -> bool:
        """슬라이드에 bullets 외 콘텐츠가 있는지 검사"""
        if not spec_slide:
            return False

        for block in spec_slide.get("content_blocks", []):
            if str(block.get("type", "")).strip().lower() != "bullets":
                return True

        normalized_blocks = normalize_slide_blocks(spec_slide) if normalize_slide_blocks else spec_slide.get("blocks", [])
        for block in normalized_blocks if isinstance(normalized_blocks, list) else []:
            if not isinstance(block, dict):
                continue
            if str(block.get("type", "")).strip().lower() not in {"bullets", "action_list"}:
                return True

        for column in spec_slide.get("columns", []):
            if self._column_has_non_bullet_content_spec(column):
                return True

        return False

    def _extract_rendered_bullet_texts(self, slide) -> List[str]:
        """
        렌더된 PPT에서 불릿으로 볼 텍스트 추출 (spec 미제공 시 폴백)
        - 제목/부제목 placeholder 제외
        - 상단 짧은 헤딩 박스와 하단 각주 박스 제외
        """
        bullet_texts: List[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            texts = [para.text.strip() for para in tf.paragraphs if para.text.strip()]
            if not texts:
                continue

            # 제목/부제목 placeholder는 제외
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                # TITLE (1), CENTER_TITLE (3), SUBTITLE (4)
                if ph_type in [1, 3, 4]:
                    continue

            shape_height = getattr(shape, "height", 0)
            shape_top = getattr(shape, "top", 0)

            # 단일 라인의 얇은 텍스트 박스는 헤딩/각주일 가능성이 높아 제외
            if len(texts) == 1:
                if shape_height and shape_height <= Pt(45):
                    continue
                if shape_top and shape_top < self.prs.slide_height * 0.22:
                    continue

            for text in texts:
                if not re.fullmatch(r"[*†‡0-9]+", text):
                    bullet_texts.append(text)

        return bullet_texts

    def _classify_text_boxes(self, slide) -> Tuple[List, List, List]:
        """
        텍스트 박스를 제목/거버닝/불릿으로 분류
        개선: 폰트 크기와 위치 기반 분류
        """
        title_boxes = []
        governing_boxes = []
        bullet_boxes = []

        title_size = self.constraints["title_font_size_pt"]
        governing_size = self.constraints["governing_font_size_pt"]
        tolerance = self.constraints["font_size_tolerance_pt"]

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            if not tf.paragraphs:
                continue

            # 첫 번째 문단의 폰트 크기로 분류
            first_para = tf.paragraphs[0]
            font_size_pt = None

            # 폰트 크기 추출
            for run in first_para.runs:
                if run.font.size:
                    font_size_pt = run.font.size.pt
                    break

            # shape.title 속성 또는 placeholder 타입으로 제목 식별
            is_title = False
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                # TITLE (1), CENTER_TITLE (3), SUBTITLE (4)
                if ph_type in [1, 3]:
                    is_title = True

            # 폰트 크기 기반 분류
            if is_title or (font_size_pt and abs(font_size_pt - title_size) <= tolerance):
                title_boxes.append(tf)
            elif font_size_pt and abs(font_size_pt - governing_size) <= tolerance:
                # 위치 기반 추가 검증: 상단 영역이면 거버닝
                if hasattr(shape, 'top'):
                    # 거버닝은 헤더 구간(상단 약 20%)으로 제한해 카드/본문 오탐을 줄임
                    slide_height = self.prs.slide_height
                    if shape.top < slide_height * 0.20:
                        governing_boxes.append(tf)
                    else:
                        bullet_boxes.append(tf)
                else:
                    governing_boxes.append(tf)
            else:
                bullet_boxes.append(tf)

        return title_boxes, governing_boxes, bullet_boxes

    def _check_bullets(self, slide_idx: int, slide, slide_constraints: dict):
        """
        불릿 검사
        개선:
        - spec가 있으면 spec 기준으로 불릿 수/길이 검증 (렌더링 오탐 방지)
        - spec가 없으면 렌더 결과에서 휴리스틱 추출
        """
        # 슬라이드별 제약 적용
        max_chars = get_max_chars_per_bullet(self.global_constraints, slide_constraints)
        min_bullets = 0
        max_bullets = get_max_bullets(self.global_constraints, slide_constraints)
        layout_name = ""
        spec_slide = self._get_spec_slide(slide_idx)

        if spec_slide:
            layout_name = normalize_layout_name(spec_slide.get("layout", "content"))
            min_bullets, max_bullets = get_bullet_bounds(
                layout_name, self.global_constraints, slide_constraints
            )

        if self.spec:
            bullet_texts = self._extract_spec_bullet_texts(slide_idx)
        else:
            bullet_texts = self._extract_rendered_bullet_texts(slide)

        for text in bullet_texts:
            if len(text) > max_chars:
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="불릿 길이",
                    message=f"불릿이 {max_chars}자를 초과합니다 ({len(text)}자)",
                    details={
                        "text_preview": text[:50] + "..." if len(text) > 50 else text,
                        "limit": max_chars
                    },
                    auto_fixable=False
                ))

            estimated_lines = max(1, (len(text) - 1) // BULLET_CHARS_PER_LINE + 1)
            if estimated_lines > BULLET_MAX_LINES:
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="불릿 줄 수",
                    message=f"불릿이 {BULLET_MAX_LINES}줄을 초과할 수 있습니다 (추정 {estimated_lines}줄)",
                    details={
                        "text_preview": text[:50] + "..." if len(text) > 50 else text,
                        "estimated_lines": estimated_lines
                    },
                    auto_fixable=False
                ))

        # 불릿 수 검사
        total_bullets = len(bullet_texts)
        if layout_name in COLUMN_LAYOUTS and spec_slide and spec_slide.get("columns"):
            # 컬럼 레이아웃은 총합이 아닌 컬럼별 과밀/공백 중심으로 점검
            per_column_limit = get_column_bullet_limit(max_bullets)
            for col_idx, col in enumerate(spec_slide.get("columns", [])):
                col_texts = self._collect_column_bullet_texts_spec(col)
                col_count = len(col_texts)
                col_has_non_bullet = self._column_has_non_bullet_content_spec(col)

                if col_count > per_column_limit:
                    self.report.add_issue(QAIssue(
                        slide_index=slide_idx,
                        severity=Severity.WARNING,
                        category="불릿 개수",
                        message=f"컬럼 불릿이 {per_column_limit}개를 초과합니다 ({col_count}개)",
                        details={
                            "column": col_idx + 1,
                            "count": col_count,
                            "max": per_column_limit,
                            "layout": layout_name,
                        },
                        auto_fixable=False
                    ))
                if col_count == 0 and not col_has_non_bullet:
                    self.report.add_issue(QAIssue(
                        slide_index=slide_idx,
                        severity=Severity.WARNING,
                        category="불릿 개수",
                        message=f"{col_idx + 1}번 컬럼에 핵심 불릿 또는 대체 콘텐츠가 없습니다",
                        details={
                            "column": col_idx + 1,
                            "layout": layout_name,
                        },
                        auto_fixable=False
                    ))
            return

        if max_bullets == 0 and total_bullets > 0:
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="불릿 개수",
                message=f"{layout_name or '현재'} 레이아웃에는 불릿이 없어야 합니다 ({total_bullets}개)",
                details={"count": total_bullets, "max": 0, "layout": layout_name or "unknown"},
                auto_fixable=False
            ))
            return

        if total_bullets > max_bullets:
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="불릿 개수",
                message=f"불릿이 {max_bullets}개를 초과합니다 ({total_bullets}개)",
                details={"count": total_bullets, "max": max_bullets, "layout": layout_name or "unknown"},
                auto_fixable=False
            ))

        has_non_bullet_content = bool(spec_slide and self._slide_has_non_bullet_content_spec(spec_slide))
        if min_bullets > 0 and total_bullets < min_bullets and not has_non_bullet_content:
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="불릿 개수",
                message=f"불릿이 {min_bullets}개 미만입니다 ({total_bullets}개)",
                details={"count": total_bullets, "min": min_bullets, "layout": layout_name or "unknown"},
                auto_fixable=False
            ))

    def _check_fonts(self, slide_idx: int, slide):
        """폰트 검사"""
        allowed_fonts = self.constraints["allowed_fonts"]
        reported_fonts = set()
        reported_sizes = set()

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    font_name = font.name

                    # 폰트 이름 검사
                    if font_name and font_name not in allowed_fonts:
                        # 유사 이름 체크 (NotoSansKR vs Noto Sans KR)
                        normalized = font_name.replace(" ", "").replace("-", "")
                        is_similar = any(
                            normalized.lower() == f.replace(" ", "").replace("-", "").lower()
                            for f in allowed_fonts
                        )
                        dedupe_key = normalized.lower()
                        if not is_similar and dedupe_key not in reported_fonts:
                            reported_fonts.add(dedupe_key)
                            self.report.add_issue(QAIssue(
                                slide_index=slide_idx,
                                severity=Severity.WARNING,
                                category="폰트",
                                message=f"허용되지 않은 폰트: {font_name}",
                                details={"font": font_name, "allowed": allowed_fonts[:3]},
                                auto_fixable=True
                            ))

                    # 폰트 사이즈 검사 (비정상적 크기)
                    if font.size:
                        size_pt = font.size.pt
                        size_key = round(size_pt, 1)
                        if (size_pt > 30 or size_pt < 8) and size_key not in reported_sizes:
                            reported_sizes.add(size_key)
                            self.report.add_issue(QAIssue(
                                slide_index=slide_idx,
                                severity=Severity.INFO,
                                category="폰트 크기",
                                message=f"비정상적인 폰트 크기: {size_pt}pt",
                                details={"size": size_pt},
                                auto_fixable=True
                            ))

        # 제목/거버닝/본문 기준 크기 준수 검사
        title_boxes, governing_boxes, bullet_boxes = self._classify_text_boxes(slide)
        tolerance = float(self.constraints["font_size_tolerance_pt"])
        title_target = float(self.constraints["title_font_size_pt"])
        governing_target = float(self.constraints["governing_font_size_pt"])
        body_target = float(self.constraints["body_font_size_pt"])

        def _collect_sizes(text_frames, min_text_len: int = 1) -> List[float]:
            sizes: List[float] = []
            for tf in text_frames:
                for para in tf.paragraphs:
                    text = para.text.strip()
                    if len(text) < min_text_len:
                        continue
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        if run.font.size:
                            sizes.append(float(run.font.size.pt))
            return sizes

        title_sizes = _collect_sizes(title_boxes, min_text_len=2)
        governing_sizes = _collect_sizes(governing_boxes, min_text_len=4)
        body_sizes = _collect_sizes(bullet_boxes, min_text_len=6)

        if title_sizes:
            median_title = sorted(title_sizes)[len(title_sizes) // 2]
            if abs(median_title - title_target) > tolerance:
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="타이포 규격",
                    message=f"제목 폰트 크기가 기준({title_target:.0f}pt)과 다를 수 있습니다 (관측 {median_title:.1f}pt)",
                    details={"expected_pt": title_target, "observed_pt": round(median_title, 1)},
                    auto_fixable=False,
                ))

        if governing_sizes:
            median_governing = sorted(governing_sizes)[len(governing_sizes) // 2]
            if abs(median_governing - governing_target) > tolerance:
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="타이포 규격",
                    message=f"거버닝 폰트 크기가 기준({governing_target:.0f}pt)과 다를 수 있습니다 (관측 {median_governing:.1f}pt)",
                    details={"expected_pt": governing_target, "observed_pt": round(median_governing, 1)},
                    auto_fixable=False,
                ))

        if body_sizes:
            median_body = sorted(body_sizes)[len(body_sizes) // 2]
            # body는 12 또는 14 허용
            if not (abs(median_body - body_target) <= tolerance or abs(median_body - (body_target + 2)) <= tolerance):
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="타이포 규격",
                    message=f"본문 폰트 크기가 기준(12~14pt) 범위를 벗어날 수 있습니다 (관측 {median_body:.1f}pt)",
                    details={"expected_pt": [body_target, body_target + 2], "observed_pt": round(median_body, 1)},
                    auto_fixable=False,
                ))

    def _check_required_blocks(self, slide_idx: int):
        """레이아웃별 필수 블록 존재 여부 검사 (spec 기반)"""
        if not self.spec:
            return
        spec_slide = self._get_spec_slide(slide_idx)
        if not spec_slide:
            return

        layout_name = normalize_layout_name(spec_slide.get("layout", "content"))
        required_blocks = []

        if self.layouts:
            layout_map = self.layouts.get("layout_map", {}) if isinstance(self.layouts.get("layout_map", {}), dict) else {}
            cfg = layout_map.get(layout_name, {}) if isinstance(layout_map.get(layout_name, {}), dict) else {}
            if isinstance(cfg.get("required_blocks"), list):
                required_blocks = [str(x).strip().lower() for x in cfg.get("required_blocks", []) if str(x).strip()]

        if not required_blocks:
            required_blocks = LAYOUT_REQUIRED_BLOCKS.get(layout_name, [])
        if not required_blocks:
            return

        available_types = set(block_types_in_slide(spec_slide) if block_types_in_slide else self._collect_spec_block_types(spec_slide))

        for req in required_blocks:
            if req == "bullets":
                has_req = bool(self._extract_spec_bullet_texts(slide_idx))
            elif req == "kpi_cards":
                has_req = bool({"kpi_cards", "kpi"} & available_types)
            elif req == "timeline_steps":
                has_req = bool({"timeline_steps", "timeline"} & available_types) or (layout_name == "timeline" and bool(self._extract_spec_bullet_texts(slide_idx)))
            else:
                has_req = req in available_types
            if not has_req:
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="레이아웃 필수 블록",
                    message=f"{layout_name} 레이아웃 필수 블록 누락: {req}",
                    details={"layout": layout_name, "required": req, "available": sorted(list(available_types))[:8]},
                    auto_fixable=False,
                ))

    def _check_block_density_rules(self, slide_idx: int):
        """blocks 기반 문장/아이템 밀도 규칙 검사"""
        if not self.spec:
            return
        spec_slide = self._get_spec_slide(slide_idx)
        if not spec_slide:
            return

        blocks = normalize_slide_blocks(spec_slide) if normalize_slide_blocks else spec_slide.get("blocks", [])
        if not isinstance(blocks, list):
            return

        cai_keywords = ("원인", "영향", "시사점")
        cai_count = 0

        for block_idx, block in enumerate(blocks):
            if not isinstance(block, dict):
                continue
            b_type = str(block.get("type", "")).strip().lower()
            if b_type not in {"bullets", "action_list"}:
                continue

            items = block.get("items", []) if isinstance(block.get("items", []), list) else []
            if not items:
                continue

            if b_type == "bullets" and (len(items) < BULLET_BLOCK_MIN_ITEMS or len(items) > BULLET_BLOCK_MAX_ITEMS):
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="블록 밀도",
                    message=f"bullets 블록 아이템 수 권장 범위({BULLET_BLOCK_MIN_ITEMS}~{BULLET_BLOCK_MAX_ITEMS}) 벗어남 ({len(items)}개)",
                    details={"block_index": block_idx, "block_type": b_type, "count": len(items)},
                    auto_fixable=False,
                ))
            if b_type == "action_list" and (len(items) < ACTION_LIST_MIN_ITEMS or len(items) > ACTION_LIST_MAX_ITEMS):
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="블록 밀도",
                    message=f"action_list 아이템 수 권장 범위({ACTION_LIST_MIN_ITEMS}~{ACTION_LIST_MAX_ITEMS}) 벗어남 ({len(items)}개)",
                    details={"block_index": block_idx, "block_type": b_type, "count": len(items)},
                    auto_fixable=False,
                ))

            for item_idx, item in enumerate(items):
                text = item if isinstance(item, str) else item.get("text", "")
                text = str(text or "").strip()
                if not text:
                    continue

                if len(text) < BULLET_RECOMMENDED_MIN_CHARS or len(text) > BULLET_RECOMMENDED_MAX_CHARS:
                    self.report.add_issue(QAIssue(
                        slide_index=slide_idx,
                        severity=Severity.INFO,
                        category="문장 규격",
                        message=f"불릿 길이 권장 범위({BULLET_RECOMMENDED_MIN_CHARS}~{BULLET_RECOMMENDED_MAX_CHARS}자) 벗어남 ({len(text)}자)",
                        details={"block_index": block_idx, "item_index": item_idx},
                        auto_fixable=False,
                    ))

                if all(k in text for k in cai_keywords):
                    cai_count += 1

        if cai_count > 1:
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.INFO,
                category="문장 구조",
                message="원인→영향→시사점 구조 불릿은 슬라이드당 1개를 권장합니다",
                details={"detected_count": cai_count},
                auto_fixable=False,
            ))

    def _collect_spec_block_types(self, spec_slide: dict) -> List[str]:
        types: List[str] = []
        if not spec_slide:
            return types
        for block in spec_slide.get("blocks", []) if isinstance(spec_slide.get("blocks", []), list) else []:
            if isinstance(block, dict):
                b_type = str(block.get("type", "")).strip().lower()
                if b_type:
                    types.append(b_type)
        for block in spec_slide.get("content_blocks", []) if isinstance(spec_slide.get("content_blocks", []), list) else []:
            if isinstance(block, dict):
                b_type = str(block.get("type", "")).strip().lower()
                if b_type:
                    types.append(b_type)
        if isinstance(spec_slide.get("bullets"), list) and spec_slide.get("bullets"):
            types.append("bullets")
        if isinstance(spec_slide.get("columns"), list) and spec_slide.get("columns"):
            types.append("columns")
        out: List[str] = []
        seen = set()
        for t in types:
            if t and t not in seen:
                seen.add(t)
                out.append(t)
        return out

    def _check_text_overflow(self, slide_idx: int, slide):
        """
        텍스트 오버플로우 추정:
        - 텍스트 길이/폰트 크기/박스 크기 기반 휴리스틱
        - 실제 PPT 엔진 줄바꿈과 100% 일치하지 않으므로 warning 수준으로 제공
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            text = "\\n".join([p.text for p in text_frame.paragraphs if p.text and p.text.strip()]).strip()
            if not text:
                continue

            width = getattr(shape, "width", None)
            height = getattr(shape, "height", None)
            if not width or not height:
                continue

            # 대표 폰트 크기 추정
            font_size_pt = float(self.constraints.get("body_font_size_pt", BODY_FONT_SIZE_PT))
            for para in text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size_pt = float(run.font.size.pt)
                        break
                if font_size_pt:
                    break

            width_pt = float(width.pt)
            height_pt = float(height.pt)
            effective_width = max(40.0, width_pt - 8.0)
            chars_per_line = max(8.0, effective_width / max(4.8, font_size_pt * 0.55))
            estimated_lines = max(1.0, len(text) / chars_per_line)
            line_height = font_size_pt * 1.32
            estimated_height = estimated_lines * line_height

            if estimated_height > (height_pt * 1.08):
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="텍스트 오버플로우",
                    message=f"텍스트가 박스 높이를 초과할 가능성이 높습니다 (추정 {int(round(estimated_height))}pt > 박스 {int(round(height_pt))}pt)",
                    details={
                        "shape_width_pt": int(round(width_pt)),
                        "shape_height_pt": int(round(height_pt)),
                        "font_size_pt": round(font_size_pt, 1),
                        "estimated_lines": int(round(estimated_lines)),
                    },
                    auto_fixable=False,
                ))

    def _check_density(self, slide_idx: int, slide):
        """콘텐츠 밀도 검사"""
        total_chars = 0
        total_paragraphs = 0
        layout_name = ""
        if self.spec and "slides" in self.spec and 0 <= slide_idx - 1 < len(self.spec["slides"]):
            layout_name = normalize_layout_name(str(self.spec["slides"][slide_idx - 1].get("layout", "")).strip().lower())

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    total_chars += len(text)
                    total_paragraphs += 1

        # 과밀도 검사
        if total_chars > DENSITY_MAX_CHARS:
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="콘텐츠 밀도",
                message=f"콘텐츠가 과밀합니다 ({total_chars}자)",
                details={"chars": total_chars, "paragraphs": total_paragraphs},
                auto_fixable=False
            ))

        # 레이아웃별 저밀도 임계치
        min_chars = DENSITY_MIN_CHARS
        min_paragraphs = DENSITY_MIN_PARAGRAPHS
        severity = Severity.INFO

        if layout_name in {"content", "comparison", "two_column", "three_column"}:
            min_chars = max(min_chars, 220)
            min_paragraphs = max(min_paragraphs, 6)
            severity = Severity.WARNING
        elif layout_name in {"chart_focus", "image_focus", "chart_insight", "competitor_2x2", "kpi_cards"}:
            min_chars = max(min_chars, 190)
            min_paragraphs = max(min_paragraphs, 5)
            severity = Severity.WARNING
        elif layout_name in {"timeline", "process_flow"}:
            min_chars = max(min_chars, 170)
            min_paragraphs = max(min_paragraphs, 4)
            severity = Severity.INFO
        elif layout_name in NO_BULLET_LAYOUTS:
            return

        if total_chars < min_chars or total_paragraphs < min_paragraphs:
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=severity,
                category="콘텐츠 밀도",
                message=f"콘텐츠 밀도가 낮습니다 ({total_chars}자 / {total_paragraphs}문단)",
                details={
                    "chars": total_chars,
                    "paragraphs": total_paragraphs,
                    "layout": layout_name or "unknown",
                    "min_chars": min_chars,
                    "min_paragraphs": min_paragraphs,
                },
                auto_fixable=False
            ))

    def _check_layout_bounds(self, slide_idx: int, slide):
        """도형이 슬라이드 경계를 벗어나는지 검사"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        tolerance = Pt(2)

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            left = getattr(shape, "left", None)
            top = getattr(shape, "top", None)
            width = getattr(shape, "width", None)
            height = getattr(shape, "height", None)

            if None in (left, top, width, height):
                continue

            right = left + width
            bottom = top + height

            if left < -tolerance or top < -tolerance or right > slide_width + tolerance or bottom > slide_height + tolerance:
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.WARNING,
                    category="레이아웃 경계",
                    message=f"도형 {shape_idx}가 슬라이드 경계를 벗어날 수 있습니다",
                    details={
                        "left": left,
                        "top": top,
                        "right": right,
                        "bottom": bottom
                    },
                    auto_fixable=False
                ))

            if shape.has_text_frame:
                text_content = " ".join(
                    para.text.strip() for para in shape.text_frame.paragraphs if para.text.strip()
                )
                if text_content and (width < Pt(60) or height < Pt(18)):
                    self.report.add_issue(QAIssue(
                        slide_index=slide_idx,
                        severity=Severity.INFO,
                        category="레이아웃 가독성",
                        message=f"도형 {shape_idx} 텍스트 영역이 너무 작을 수 있습니다",
                        details={"text_preview": text_content[:40]},
                        auto_fixable=False
                    ))

    def _check_forbidden_words(self, slide_idx: int, slide, slide_constraints: dict):
        """금지어 검사 (global + slide_constraints 병합)"""
        forbidden = get_forbidden_words(self.global_constraints, slide_constraints)
        if not forbidden:
            return

        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_text += para.text + " "

        all_text_lower = all_text.lower()

        for word in forbidden:
            if word.lower() in all_text_lower:
                self.report.add_issue(QAIssue(
                    slide_index=slide_idx,
                    severity=Severity.ERROR,
                    category="금지어",
                    message=f"금지어 발견: '{word}'",
                    details={"word": word},
                    auto_fixable=True
                ))

    def _check_spec_alignment(self, slide_idx: int, slide):
        """
        Spec과의 일치 검사
        개선: 폰트 크기/위치 기준으로 제목 식별 고도화
        """
        if not self.spec or "slides" not in self.spec:
            return

        slides = self.spec["slides"]
        if slide_idx - 1 >= len(slides):
            return

        spec_slide = slides[slide_idx - 1]
        spec_title = spec_slide.get("title", "")

        # 텍스트 박스 분류로 제목 식별
        title_boxes, _, _ = self._classify_text_boxes(slide)

        actual_title = None
        for tf in title_boxes:
            if tf.paragraphs:
                text = tf.paragraphs[0].text.strip()
                if text:
                    actual_title = text
                    break

        # 폴백: 첫 번째 텍스트 shape에서 추출
        if not actual_title:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                    if text:
                        actual_title = text
                        break

        if actual_title and spec_title:
            # 정규화 비교
            norm_actual = re.sub(r'\s+', '', actual_title.lower())
            norm_spec = re.sub(r'\s+', '', spec_title.lower())

            # 부분 일치 허용 (포함 여부)
            if norm_spec not in norm_actual and norm_actual not in norm_spec:
                # 유사도 검사 (간단한 Jaccard)
                actual_words = set(actual_title.lower().split())
                spec_words = set(spec_title.lower().split())
                if actual_words and spec_words:
                    intersection = len(actual_words & spec_words)
                    union = len(actual_words | spec_words)
                    similarity = intersection / union if union > 0 else 0

                    if similarity < 0.5:  # 50% 미만 유사도
                        self.report.add_issue(QAIssue(
                            slide_index=slide_idx,
                            severity=Severity.INFO,
                            category="Spec 불일치",
                            message="제목이 Spec과 다를 수 있습니다",
                            details={
                                "spec_title": spec_title,
                                "actual_title": actual_title[:50] if len(actual_title) > 50 else actual_title,
                                "similarity": f"{similarity:.0%}"
                            },
                            auto_fixable=False
                        ))

    def _check_global(self):
        """전역 검사"""
        # 슬라이드 수 검사
        if self.global_constraints:
            max_slides = self.global_constraints.get("max_slides", 50)
            if len(self.prs.slides) > max_slides:
                self.report.add_issue(QAIssue(
                    slide_index=0,
                    severity=Severity.WARNING,
                    category="슬라이드 수",
                    message=f"슬라이드가 {max_slides}장을 초과합니다 ({len(self.prs.slides)}장)",
                    details={"count": len(self.prs.slides), "max": max_slides},
                    auto_fixable=False
                ))

            # required_sections 검사 (layout 또는 metadata.section 기준)
            required_sections = self.global_constraints.get("required_sections", [])
            if self.spec and isinstance(required_sections, list) and required_sections:
                slides = self.spec.get("slides", [])
                available_layouts = {normalize_layout_name(str(slide.get("layout", "")).strip().lower()) for slide in slides}
                available_sections = set()
                for slide in slides:
                    metadata = slide.get("metadata", {})
                    section_name = metadata.get("section")
                    if section_name:
                        available_sections.add(str(section_name).strip().lower())

                for section in required_sections:
                    section_norm = str(section).strip().lower()
                    if not section_norm:
                        continue
                    if section_norm not in available_layouts and section_norm not in available_sections:
                        self.report.add_issue(QAIssue(
                            slide_index=0,
                            severity=Severity.WARNING,
                            category="필수 섹션",
                            message=f"필수 섹션 미존재: '{section}'",
                            details={"required": section},
                            auto_fixable=False
                        ))

    def _check_evidence(self):
        """
        Evidence 검사
        개선: sources.md# 앵커 포맷 검증 및 존재 검사
        """
        if not self.spec or "slides" not in self.spec:
            return

        evidence_pattern = re.compile(EVIDENCE_ANCHOR_PATTERN)

        # top-level sources_ref 검사
        for ref in self.spec.get("sources_ref", []):
            self._validate_evidence_anchor(0, ref, evidence_pattern)

        for slide_idx, spec_slide in enumerate(self.spec["slides"], start=1):
            # 슬라이드 레벨 metadata의 source_refs 검사
            metadata = spec_slide.get("metadata", {})
            source_refs = metadata.get("source_refs", [])
            for ref in source_refs:
                self._validate_evidence_anchor(slide_idx, ref, evidence_pattern)

            # bullets 내 evidence 검사
            self._check_bullets_evidence(slide_idx, spec_slide.get("bullets", []), evidence_pattern)

            # columns 내 bullets evidence 검사
            for col in spec_slide.get("columns", []):
                self._check_bullets_evidence(slide_idx, col.get("bullets", []), evidence_pattern)
                self._check_visual_evidence(slide_idx, col.get("visual"), evidence_pattern)

                # columns[].content_blocks 내 evidence 검사
                for block in col.get("content_blocks", []):
                    self._check_content_block_evidence(slide_idx, block, evidence_pattern)

            # slide visuals evidence 검사
            for visual in spec_slide.get("visuals", []):
                self._check_visual_evidence(slide_idx, visual, evidence_pattern)

            # content_blocks 내 evidence 검사
            for block in spec_slide.get("content_blocks", []):
                self._check_content_block_evidence(slide_idx, block, evidence_pattern)

            # blocks 내 evidence 검사
            normalized_blocks = normalize_slide_blocks(spec_slide) if normalize_slide_blocks else spec_slide.get("blocks", [])
            for block in normalized_blocks if isinstance(normalized_blocks, list) else []:
                if not isinstance(block, dict):
                    continue
                if "evidence" in block:
                    self._validate_evidence(slide_idx, block["evidence"], evidence_pattern)

                for item in block.get("items", []) if isinstance(block.get("items", []), list) else []:
                    if isinstance(item, dict) and "evidence" in item:
                        self._validate_evidence(slide_idx, item["evidence"], evidence_pattern)

                for card in block.get("cards", []) if isinstance(block.get("cards", []), list) else []:
                    if isinstance(card, dict) and "evidence" in card:
                        self._validate_evidence(slide_idx, card["evidence"], evidence_pattern)

                if isinstance(block.get("chart"), dict):
                    self._check_visual_evidence(slide_idx, block.get("chart"), evidence_pattern)
                if isinstance(block.get("image"), dict):
                    self._check_visual_evidence(slide_idx, block.get("image"), evidence_pattern)

            # footnotes 내 evidence 검사
            for footnote in spec_slide.get("footnotes", []):
                if isinstance(footnote, dict) and "evidence" in footnote:
                    self._validate_evidence(slide_idx, footnote["evidence"], evidence_pattern)

    def _check_bullets_evidence(self, slide_idx: int, bullets: list, pattern):
        """불릿 내 evidence 검사"""
        for bullet in bullets:
            if isinstance(bullet, dict) and "evidence" in bullet:
                self._validate_evidence(slide_idx, bullet["evidence"], pattern)

    def _check_visual_evidence(self, slide_idx: int, visual: dict, pattern):
        """visual 내 evidence 검사"""
        if isinstance(visual, dict) and "evidence" in visual:
            self._validate_evidence(slide_idx, visual["evidence"], pattern)

    def _check_content_block_evidence(self, slide_idx: int, block: dict, pattern):
        """content_block 내부 evidence 검사"""
        if not isinstance(block, dict):
            return

        # 블록 레벨 evidence
        if "evidence" in block:
            self._validate_evidence(slide_idx, block["evidence"], pattern)

        # 블록 내 bullets
        self._check_bullets_evidence(slide_idx, block.get("bullets", []), pattern)

        # nested object evidence
        table_def = block.get("table")
        if isinstance(table_def, dict) and "evidence" in table_def:
            self._validate_evidence(slide_idx, table_def["evidence"], pattern)

        self._check_visual_evidence(slide_idx, block.get("chart"), pattern)
        self._check_visual_evidence(slide_idx, block.get("image"), pattern)

        quote_def = block.get("quote")
        if isinstance(quote_def, dict) and "evidence" in quote_def:
            self._validate_evidence(slide_idx, quote_def["evidence"], pattern)

        kpi_def = block.get("kpi")
        if isinstance(kpi_def, dict) and "evidence" in kpi_def:
            self._validate_evidence(slide_idx, kpi_def["evidence"], pattern)

    def _validate_evidence(self, slide_idx: int, evidence: dict, pattern):
        """evidence 객체 검증"""
        if not isinstance(evidence, dict):
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="Evidence 포맷",
                message="evidence는 object 형태여야 합니다",
                details={"evidence": str(evidence)},
                auto_fixable=False
            ))
            return
        if "source_anchor" in evidence:
            self._validate_evidence_anchor(slide_idx, evidence["source_anchor"], pattern)

    def _validate_evidence_anchor(self, slide_idx: int, anchor: str, pattern):
        """앵커 포맷 및 존재 검사"""
        if not isinstance(anchor, str):
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="Evidence 포맷",
                message="앵커 타입이 문자열이 아닙니다",
                details={"anchor": str(anchor), "expected_format": "sources.md#anchor-name"},
                auto_fixable=False
            ))
            return

        # 포맷 검사
        if not pattern.match(anchor):
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="Evidence 포맷",
                message=f"잘못된 앵커 포맷: '{anchor}'",
                details={
                    "anchor": anchor,
                    "expected_format": "sources.md#anchor-name"
                },
                auto_fixable=False
            ))
            return

        # 존재 검사 (sources.md가 로드된 경우)
        if self.sources_anchors and anchor not in self.sources_anchors:
            self.report.add_issue(QAIssue(
                slide_index=slide_idx,
                severity=Severity.INFO,
                category="Evidence 참조",
                message=f"앵커를 찾을 수 없음: '{anchor}'",
                details={
                    "anchor": anchor,
                    "available_anchors": list(self.sources_anchors)[:5]
                },
                auto_fixable=False
            ))


def validate_spec_business_rules(spec: dict, global_constraints: dict = None) -> List[QAIssue]:
    """
    Spec 비즈니스 규칙 검증
    개선: columns[].bullets도 포함
    """
    issues = []
    gc = global_constraints or spec.get("global_constraints", {})

    for slide_idx, slide in enumerate(spec.get("slides", []), start=1):
        layout = normalize_layout_name(slide.get("layout", "content"))
        slide_constraints = slide.get("slide_constraints", {})

        # 슬라이드별 제약 적용
        max_bullets = get_max_bullets(gc, slide_constraints)
        max_chars = get_max_chars_per_bullet(gc, slide_constraints)

        # 일반 bullets 검사
        bullets = slide.get("bullets", [])
        if bullets:
            _validate_bullets_list(issues, slide_idx, bullets, max_bullets, max_chars, "bullets")

        # columns 내 bullets 검사
        for col_idx, col in enumerate(slide.get("columns", []), start=1):
            col_bullets = col.get("bullets", [])
            if col_bullets:
                _validate_bullets_list(
                    issues, slide_idx, col_bullets, max_bullets, max_chars,
                    f"columns[{col_idx}].bullets"
                )

        # content_blocks 내 bullets 검사
        for block_idx, block in enumerate(slide.get("content_blocks", []), start=1):
            if block.get("type") == "bullets":
                block_bullets = block.get("bullets", [])
                if block_bullets:
                    _validate_bullets_list(
                        issues, slide_idx, block_bullets, max_bullets, max_chars,
                        f"content_blocks[{block_idx}].bullets"
                    )

        # blocks 내 bullets/action_list 검사
        normalized_blocks = normalize_slide_blocks(slide) if normalize_slide_blocks else slide.get("blocks", [])
        for block_idx, block in enumerate(normalized_blocks if isinstance(normalized_blocks, list) else [], start=1):
            if not isinstance(block, dict):
                continue
            b_type = str(block.get("type", "")).strip().lower()
            if b_type in {"bullets", "action_list"}:
                block_items = block.get("items", [])
                if block_items:
                    _validate_bullets_list(
                        issues, slide_idx, block_items, max_bullets, max_chars,
                        f"blocks[{block_idx}].items"
                    )

    return issues


def _validate_bullets_list(issues: List[QAIssue], slide_idx: int, bullets: list,
                           max_bullets: int, max_chars: int, location: str):
    """불릿 리스트 검증 헬퍼"""
    # 개수 검사
    if len(bullets) > max_bullets:
        issues.append(QAIssue(
            slide_index=slide_idx,
            severity=Severity.WARNING,
            category="불릿 개수 (Spec)",
            message=f"{location}: 불릿이 {max_bullets}개를 초과합니다 ({len(bullets)}개)",
            details={"count": len(bullets), "max": max_bullets, "location": location},
            auto_fixable=False
        ))

    # 길이 검사
    for i, bullet in enumerate(bullets):
        text = bullet if isinstance(bullet, str) else bullet.get("text", "")
        if len(text) > max_chars:
            issues.append(QAIssue(
                slide_index=slide_idx,
                severity=Severity.WARNING,
                category="불릿 길이 (Spec)",
                message=f"{location}[{i}]: 불릿이 {max_chars}자를 초과합니다 ({len(text)}자)",
                details={
                    "text_preview": text[:50] + "..." if len(text) > 50 else text,
                    "length": len(text),
                    "max": max_chars
                },
                auto_fixable=False
            ))


def main():
    parser = argparse.ArgumentParser(
        description="PPT QA 자동 검사기 v2.1",
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument("pptx_path", help="검사할 PPTX 파일 경로")
    parser.add_argument("--spec", help="deck_spec.yaml 경로 (선택)")
    parser.add_argument("--tokens", help="tokens.yaml 경로 (선택)")
    parser.add_argument("--layouts", help="layouts.yaml 경로 (선택)")
    parser.add_argument("--sources", help="sources.md 경로 (Evidence 검증용, 선택)")
    parser.add_argument("--output", "-o", help="보고서 출력 경로 (JSON)")
    parser.add_argument("--markdown", "-m", help="마크다운 보고서 출력 경로")
    parser.add_argument("--fix", action="store_true", help="자동 수정 가능한 이슈 수정 (미구현)")
    parser.add_argument("--verbose", "-v", action="store_true", help="상세 출력")

    args = parser.parse_args()

    # 파일 존재 확인
    if not Path(args.pptx_path).exists():
        print(f"❌ 파일을 찾을 수 없습니다: {args.pptx_path}")
        return 1

    # QA 실행
    checker = PPTQAChecker(
        pptx_path=args.pptx_path,
        spec_path=args.spec,
        tokens_path=args.tokens,
        sources_path=args.sources,
        layouts_path=args.layouts
    )

    report = checker.run_all_checks()

    # 결과 출력
    if args.verbose or not args.output:
        print(report.to_markdown())

    # JSON 출력
    if args.output:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(report.to_dict(), f, ensure_ascii=False, indent=2)
        print(f"\n📄 JSON 보고서 저장: {args.output}")

    # 마크다운 출력
    if args.markdown:
        md_path = Path(args.markdown)
        md_path.parent.mkdir(parents=True, exist_ok=True)
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(report.to_markdown())
        print(f"📝 마크다운 보고서 저장: {args.markdown}")

    # 결과 반환
    if report.passed:
        print(f"\n✅ QA 통과 (경고: {report.warning_count}, 참고: {report.info_count})")
        return 0
    else:
        print(f"\n❌ QA 실패 (오류: {report.error_count}, 경고: {report.warning_count})")
        return 1


if __name__ == "__main__":
    sys.exit(main())
